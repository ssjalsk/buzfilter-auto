# -*- coding: utf-8 -*-
import streamlit as st
import pandas as pd
import gspread
from oauth2client.service_account import ServiceAccountCredentials
import anthropic
from datetime import datetime
import os
import json
import re
import io
import requests
import zipfile
import base64
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

def get_anthropic_client():
    try:
        return anthropic.Anthropic(api_key=st.secrets["ANTHROPIC_API_KEY"])
    except:
        return anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))

def get_sheet(worksheet_name):
    SHEET_URL = "https://docs.google.com/spreadsheets/d/1CtD6VVtmiQNz90mKJFfuPq8-LMowLHg3NZPnoqwpISE/"
    try:
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
        try:
            creds_dict = json.loads(st.secrets["GOOGLE_CREDENTIALS"])
            creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
        except:
            BASE_DIR = os.path.dirname(os.path.abspath(__file__))
            creds = ServiceAccountCredentials.from_json_keyfile_name(os.path.join(BASE_DIR, 'service_account.json'), scope)
        client_gs = gspread.authorize(creds)
        return client_gs.open_by_url(SHEET_URL).worksheet(worksheet_name)
    except Exception as e:
        st.error(f"시트 연결 실패 ({worksheet_name}): {e}")
        return None

def find_last_data_row(sheet):
    all_values = sheet.get_all_values()
    last_row = 2
    for i, row in enumerate(all_values):
        if len(row) > 1 and row[1].strip() != '':
            last_row = i + 1
    return last_row + 1

def insert_row_safe(sheet, start_row, rows_data):
    if not rows_data:
        return
    col_map = [('B', 0), ('C', 1), ('D', 2), ('E', 3), ('F', 4), ('H', 5), ('I', 6), ('K', 7)]
    updates = []
    for i, row in enumerate(rows_data):
        r = start_row + i
        for col, idx in col_map:
            updates.append({'range': f'{col}{r}', 'values': [[row[idx]]]})
    sheet.batch_update(updates, value_input_option='USER_ENTERED')

def extract_qty_from_text(text):
    text = str(text)
    m = re.search(r'/\s*(\d+)\s*(세트|개|박스|팩|장|묶음)', text)
    if m:
        return int(m.group(1))
    m = re.search(r'(\d+)\s*(세트|박스|팩|묶음)', text)
    if m:
        return int(m.group(1))
    return 1

def normalize_for_match(text):
    """매칭용 정규화: 소문자 + 공백·특수문자 제거 (띄어쓰기 유무 무관하게 비교)"""
    text = str(text).lower()
    return re.sub(r'[\s\-_·,.]', '', text)

def split_mixed_tokens(text):
    """
    언어 경계 분할로 혼합 토큰 추출
    예) '쿠쿠X툴' → ['쿠쿠', 'X', '툴', 'X툴', '쿠쿠X툴']
    예) '삼성3벌' → ['삼성', '3', '벌', '3벌', '삼성3벌']
    → 띄어쓰기 유무 무관하게 동일한 토큰 집합 생성
    """
    # 1) 언어·숫자 단위로 분할
    atoms = re.findall(r'[가-힣]+|[a-zA-Z]+|\d+', text)
    tokens = set(atoms)
    # 2) 인접 원자들의 조합도 추가 (2~3개 연속)
    for i in range(len(atoms)):
        for j in range(i + 1, min(i + 4, len(atoms) + 1)):
            combined = ''.join(atoms[i:j])
            if len(combined) >= 2:
                tokens.add(combined)
    return [normalize_for_match(t) for t in tokens if len(t) >= 2]


def find_top_candidates(raw, calc_df, product_col, top_n=5):
    """
    3단계 매칭:
    1단계 - 브랜드 하드 필터 (브랜드명 첫줄만 비교, \n 포함 브랜드 대응)
    2단계 - 필터 시리즈/제품명 키워드 점수화 + 경쟁 변종 패널티
    3단계 - 명확한 1위 존재 시 AI 없이 자동 확정, 아니면 top_n 반환
    """
    query_no_qty = re.split(r'\s*/\s*\d+', raw)[0]
    full_query_norm = normalize_for_match(query_no_qty)
    query_tokens = split_mixed_tokens(query_no_qty)

    # 브랜드명에 '\n' 있는 경우(예: '쿠쿠\n공기청정기 필터') 첫 줄만 사용
    def brand_short(b):
        return normalize_for_match(str(b).split('\n')[0])

    # 1단계: 브랜드 하드 필터 — 가장 긴 매칭 브랜드 우선
    matched_brand = None
    best_len = 0
    for brand in calc_df['브랜드'].dropna().unique():
        bn = brand_short(brand)
        if len(bn) >= 2 and bn in full_query_norm and len(bn) > best_len:
            matched_brand = brand
            best_len = len(bn)

    if matched_brand:
        mb_short = brand_short(matched_brand)
        search_df = calc_df[calc_df['브랜드'].apply(lambda x: brand_short(x) == mb_short)]
        if len(search_df) == 0:
            search_df = calc_df
    else:
        search_df = calc_df

    if len(search_df) == 1:
        return search_df

    # '필터 시리즈' 컬럼 감지 (모델 구분자 핵심 컬럼)
    series_col = '필터 시리즈' if '필터 시리즈' in calc_df.columns else ''

    # 2단계: 키워드 점수화
    scored = []
    for idx, row in search_df.iterrows():
        series_n = normalize_for_match(str(row.get(series_col, ''))) if series_col else ''
        product_n = normalize_for_match(str(row.get('제품명', '')))
        code_n = normalize_for_match(str(row.get(product_col, '')))
        combined_n = series_n + '|' + product_n

        score = 0
        for tok in query_tokens:
            if len(tok) < 2:
                continue
            if tok in series_n:
                score += 10 + len(tok) * 4  # 시리즈 가중치 높음
            elif tok in product_n or tok in code_n:
                score += 5 + len(tok) * 3

        # 경쟁 변종 패널티: X툴≠Y툴, 3벌≠5벌 등
        for q_tok in query_tokens:
            m = re.match(r'^([a-z]+)([가-힣]{1,4})$', q_tok)
            if m:
                kr_suffix = m.group(2)
                for competitor in re.findall(r'[a-z]+' + kr_suffix, combined_n):
                    if competitor != q_tok:
                        score -= 40
            m2 = re.match(r'^(\d+)([가-힣]{1,4})$', q_tok)
            if m2:
                kr_suffix2 = m2.group(2)
                for competitor in re.findall(r'\d+' + kr_suffix2, combined_n):
                    if competitor != q_tok:
                        score -= 40

        if score > 0:
            scored.append((score, idx))

    scored.sort(reverse=True)

    if not scored:
        return search_df.head(top_n)

    # 3단계: 명확한 1위 자동 확정
    if len(scored) == 1:
        return search_df.loc[[scored[0][1]]]
    if scored[0][0] >= 20 and scored[0][0] >= scored[1][0] * 1.7:
        return search_df.loc[[scored[0][1]]]

    top_idx = [idx for _, idx in scored[:top_n]]
    return search_df.loc[top_idx]


def generate_quote_pdf(quote_data, stamp_path=None):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.units import mm
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    fr = os.path.join(BASE_DIR, 'NotoSansKR-Regular.ttf')
    fb_path = os.path.join(BASE_DIR, 'NotoSansKR-Bold.ttf')
    has_kor = os.path.exists(fr) and os.path.exists(fb_path)
    if has_kor:
        pdfmetrics.registerFont(TTFont('KR', fr))
        pdfmetrics.registerFont(TTFont('KR-B', fb_path))
        fn, fb = 'KR', 'KR-B'
    else:
        fn, fb = 'Helvetica', 'Helvetica-Bold'

    buf = io.BytesIO()
    w, h = A4
    c = canvas.Canvas(buf, pagesize=A4)
    LG = colors.HexColor("#F2F2F2")
    MG = colors.HexColor("#CCCCCC")
    DG = colors.HexColor("#404040")
    TBG = colors.HexColor("#D9D9D9")
    ML, MR = 20*mm, w-20*mm
    PW = MR - ML
    items = quote_data["items"]
    is_tax = quote_data["tax_type"] == "발행"
    sup = sum(int(it["수량"])*int(it["단가"]) for it in items)
    vat = int(sup*0.1) if is_tax else 0
    tot = sup + vat
    y = h - 18*mm
    c.setFont(fb, 28); c.setFillColor(colors.black)
    c.drawCentredString(w/2, y, "견   적   서")
    y -= 10*mm
    c.setStrokeColor(colors.black); c.setLineWidth(1.5); c.line(ML, y, MR, y)
    y -= 8*mm
    bt = y
    rcx, rcw = ML+PW*0.5, PW*0.5
    c.setFont(fb, 18); c.setFillColor(colors.HexColor("#1a5fa8"))
    c.drawString(ML+5*mm, bt-12*mm, "Aligo")
    c.setFont(fb, 14); c.drawString(ML+5*mm, bt-20*mm, "Media")
    if stamp_path and os.path.exists(stamp_path):
        try: c.drawImage(stamp_path, rcx-24*mm, bt-33*mm, width=22*mm, height=22*mm, mask='auto')
        except: pass
    c.setFillColor(DG); c.rect(rcx, bt-6*mm, rcw, 6*mm, fill=1, stroke=0)
    c.setFillColor(colors.white); c.setFont(fb, 10)
    c.drawCentredString(rcx+rcw/2, bt-4.5*mm, "공  급  자")
    srows = [("등록번호","161-22-02310","대표자","박철규"),("상  호","알리고미디어","",""),
             ("주  소","서울 마포구 양화로64, 8층","",""),("연락처","010-9469-2381","",""),
             ("업  태","전문, 서비스업","종 목","광고대행업")]
    rh = 5.5*mm
    for i,(k1,v1,k2,v2) in enumerate(srows):
        ry = bt-6*mm-(i+1)*rh
        c.setFillColor(LG if i%2==0 else colors.white); c.rect(rcx, ry, rcw, rh, fill=1, stroke=0)
        c.setStrokeColor(MG); c.setLineWidth(0.5); c.rect(rcx, ry, rcw, rh, fill=0, stroke=1)
        c.setFillColor(colors.black)
        c.setFont(fb,8); c.drawString(rcx+2*mm, ry+1.5*mm, k1)
        c.setFont(fn,8); c.drawString(rcx+16*mm, ry+1.5*mm, v1)
        if k2:
            c.setFont(fb,8); c.drawString(rcx+rcw*0.62, ry+1.5*mm, k2)
            c.setFont(fn,8); c.drawString(rcx+rcw*0.62+10*mm, ry+1.5*mm, v2)
    y = bt-6*mm-len(srows)*rh-5*mm
    c.setStrokeColor(MG); c.setLineWidth(0.5)
    c.setFillColor(LG); c.rect(ML, y-6*mm, 28*mm, 6*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,9)
    c.drawCentredString(ML+14*mm, y-4.5*mm, "견  적  일")
    c.setFillColor(colors.white); c.rect(ML+28*mm, y-6*mm, PW/2-28*mm, 6*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fn,9); c.drawString(ML+30*mm, y-4.5*mm, quote_data["date"])
    ax = ML+PW/2
    c.setFillColor(colors.white); c.rect(ax, y-6*mm, PW/2, 6*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,8)
    c.drawCentredString(ax+PW/4, y-4.5*mm, "IBK기업은행")
    y -= 6*mm
    c.setFillColor(colors.white); c.rect(ML, y-6*mm, PW/2, 6*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,10)
    c.drawCentredString(ML+PW/4, y-4.5*mm, f"{quote_data['client']}  귀하")
    ax2 = ML+PW/2
    c.setFillColor(colors.white); c.rect(ax2, y-6*mm, PW/2, 6*mm, fill=1, stroke=1)
    c.setFont(fn,7.5); c.setFillColor(colors.black)
    c.drawCentredString(ax2+PW/4, y-4.5*mm, "208-174145-04-018 박철규 (알리고 미디어)")
    y -= 8*mm
    c.setFillColor(TBG); c.rect(ML, y-12*mm, PW*0.38, 12*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,11)
    c.drawCentredString(ML+PW*0.19, y-6*mm, "합계금액")
    c.setFont(fn,8); c.drawCentredString(ML+PW*0.19, y-10*mm, "(부가세 포함)" if is_tax else "(VAT 미포함)")
    c.setFillColor(LG); c.rect(ML+PW*0.38, y-12*mm, PW*0.47, 12*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,11)
    c.drawCentredString(ML+PW*0.615, y-7*mm, "진행 상품 상세 내역")
    c.setFillColor(colors.white); c.rect(ML+PW*0.85, y-12*mm, PW*0.15, 12*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,9)
    c.drawCentredString(ML+PW*0.925, y-7*mm, f"₩{tot:,}")
    y -= 14*mm
    cx = [ML, ML+12*mm, ML+70*mm, ML+105*mm, ML+118*mm, ML+136*mm, ML+154*mm]
    cw = [12*mm,58*mm,35*mm,13*mm,18*mm,18*mm]; cw.append(MR-cx[-1])
    labels = ["NO","품목","구성","수량","단가","공급가액(VAT별도)","비고"]
    c.setFillColor(TBG); c.rect(ML, y-6*mm, PW, 6*mm, fill=1, stroke=1)
    c.setFillColor(colors.black); c.setFont(fb,8)
    for i,lbl in enumerate(labels):
        c.drawCentredString(cx[i]+cw[i]/2, y-4.5*mm, lbl)
        if i>0: c.setLineWidth(0.5); c.line(cx[i],y,cx[i],y-6*mm)
    y -= 6*mm
    max_r = max(len(items),10); rh2=6*mm
    for i in range(max_r):
        c.setFillColor(LG if i%2==0 else colors.white); c.rect(ML, y-rh2, PW, rh2, fill=1, stroke=0)
        c.setStrokeColor(MG); c.setLineWidth(0.3); c.rect(ML, y-rh2, PW, rh2, fill=0, stroke=1)
        c.setFillColor(colors.black); c.setFont(fn,8)
        c.drawCentredString(cx[0]+cw[0]/2, y-4.5*mm, str(i+1))
        if i < len(items):
            it=items[i]; qty=int(it["수량"]); price=int(it["단가"]); sp=qty*price
            c.drawString(cx[1]+2*mm, y-4.5*mm, str(it.get("품목","")))
            c.drawString(cx[2]+2*mm, y-4.5*mm, str(it.get("구성","")))
            c.drawCentredString(cx[3]+cw[3]/2, y-4.5*mm, f"{qty:,}")
            c.drawRightString(cx[4]+cw[4]-1*mm, y-4.5*mm, f"{price:,}")
            c.drawRightString(cx[5]+cw[5]-1*mm, y-4.5*mm, f"{sp:,}")
            c.drawCentredString(cx[6]+cw[6]/2, y-4.5*mm, str(it.get("비고","")))
        else:
            c.drawRightString(cx[5]+cw[5]-1*mm, y-4.5*mm, "0")
        for j in range(1,len(cx)): c.setLineWidth(0.3); c.line(cx[j],y,cx[j],y-rh2)
        y -= rh2
    sums = [("공급가액 합계",sup),("세  액 (VAT)",vat),("합  계(부가세 포함)",tot)] if is_tax else [("공급가액 합계 (VAT 미발행)",sup)]
    for lbl,amt in sums:
        c.setFillColor(TBG); c.rect(ML, y-7*mm, PW, 7*mm, fill=1, stroke=1)
        c.setFillColor(colors.black); c.setFont(fb,9)
        c.drawCentredString(ML+PW*0.5, y-4.8*mm, lbl)
        c.drawRightString(cx[5]+cw[5]-1*mm, y-4.8*mm, f"{amt:,}")
        y -= 7*mm
    y -= 5*mm
    c.setFont(fb,9); c.setFillColor(colors.black)
    c.drawString(ML, y, "▶ 입금 계좌번호 : IBK기업은행 208-174145-04-018 박철규 (알리고 미디어)")
    y -= 6*mm
    c.drawString(ML, y, f"▶ 비  고 : {quote_data.get('memo','')}")
    c.save(); buf.seek(0)
    return buf

def deploy_to_netlify(html_content, site_id, token, extra_files=None):
    try:
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.writestr('index.html', html_content.encode('utf-8'))
            if extra_files:
                for filename, file_bytes in extra_files.items():
                    zf.writestr(filename, file_bytes)
        zip_buffer.seek(0)
        headers = {
            'Authorization': f'Bearer {token}',
            'Content-Type': 'application/zip'
        }
        response = requests.post(
            f'https://api.netlify.com/api/v1/sites/{site_id}/deploys',
            headers=headers,
            data=zip_buffer.getvalue(),
            timeout=60
        )
        if response.status_code in [200, 201]:
            return True, "성공"
        else:
            return False, f"오류 코드: {response.status_code}\n{response.text[:300]}"
    except Exception as e:
        return False, str(e)

# ─────────────────────────────────────────────
# 📖 바이럴 백과사전 헬퍼 함수
# ─────────────────────────────────────────────
VIRAL_SHEET_ID = "1CtD6VVtmiQNz90mKJFfuPq8-LMowLHg3NZPnoqwpISE"  # 버즈필터 시트 (서비스 계정 권한 있음)
VIRAL_TAB_NAME = "바이럴백과사전"

def extract_first_image(html_content):
    """HTML 본문에서 첫 번째 이미지 URL 추출 (썸네일용)"""
    m = re.search(r'<img[^>]+src=["\']([^"\']+)["\']', html_content or "")
    return m.group(1) if m else None

def get_viral_sheet():
    """바이럴 백과사전 구글시트 탭 연결. 없으면 자동 생성."""
    try:
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
        try:
            creds_dict = json.loads(st.secrets["GOOGLE_CREDENTIALS"])
            creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
        except Exception:
            BASE_DIR = os.path.dirname(os.path.abspath(__file__))
            creds = ServiceAccountCredentials.from_json_keyfile_name(
                os.path.join(BASE_DIR, 'service_account.json'), scope)
        client_gs = gspread.authorize(creds)
        spreadsheet = client_gs.open_by_url(
            f"https://docs.google.com/spreadsheets/d/{VIRAL_SHEET_ID}/")
        try:
            return spreadsheet.worksheet(VIRAL_TAB_NAME)
        except gspread.exceptions.WorksheetNotFound:
            ws = spreadsheet.add_worksheet(title=VIRAL_TAB_NAME, rows=500, cols=10)
            ws.append_row(["날짜", "slug", "제목", "해시태그", "요약", "본문HTML"])
            return ws
    except Exception as e:
        st.error(f"바이럴 백과사전 시트 연결 실패: {e}")
        return None


def get_viral_posts():
    """게시글 목록 반환 (최신 순)"""
    ws = get_viral_sheet()
    if not ws:
        return []
    rows = ws.get_all_values()
    if len(rows) <= 1:
        return []
    posts = []
    for row in rows[1:]:
        if len(row) >= 3 and row[2].strip():
            posts.append({
                "날짜": row[0] if len(row) > 0 else "",
                "slug": row[1] if len(row) > 1 else "",
                "제목": row[2] if len(row) > 2 else "",
                "해시태그": row[3] if len(row) > 3 else "",
                "요약": row[4] if len(row) > 4 else "",
                "본문HTML": row[5] if len(row) > 5 else "",
            })
    return list(reversed(posts))


def save_viral_post(slug, title, hashtags, summary, html_body):
    """새 게시글 구글시트에 저장"""
    ws = get_viral_sheet()
    if not ws:
        return False
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    ws.append_row([now, slug, title, hashtags, summary, html_body])
    return True


def delete_viral_post(slug):
    """slug로 게시글 구글시트에서 삭제"""
    ws = get_viral_sheet()
    if not ws:
        return False
    rows = ws.get_all_values()
    for i, row in enumerate(rows):
        if len(row) > 1 and row[1] == slug:
            ws.delete_rows(i + 1)
            return True
    return False


def make_slug(title):
    """제목에서 URL slug 생성 (타임스탬프 + 제목)"""
    ts = datetime.now().strftime("%Y%m%d%H%M%S")
    slug_title = re.sub(r'[^\w가-힣]', '-', title).strip('-')
    slug_title = re.sub(r'-+', '-', slug_title)[:40]
    return f"{ts}-{slug_title}"


@st.cache_data(ttl=300)
def get_view_count(slug):
    """hits.seeyoufarm.com에서 조회수 가져오기 (5분 캐시)"""
    import urllib.parse as _ulp
    try:
        encoded = _ulp.quote(
            f"https://aligomedia.co.kr/blog/{slug}/", safe="")
        r = requests.get(
            f"https://hits.seeyoufarm.com/api/count/incr/badge.svg?url={encoded}",
            timeout=8)
        if r.status_code == 200:
            nums = re.findall(r">(\d+)<", r.text)
            valid = [int(n) for n in nums if n.isdigit() and len(n) <= 7]
            return max(valid) if valid else 0
        return 0
    except Exception:
        return 0


def generate_post_html(post):
    """개별 포스트 전체 HTML 생성 (Naver SEO 최적화)"""
    tags = [t.strip() for t in post["해시태그"].split(",") if t.strip()]
    tag_meta = ", ".join(tags)
    tag_header_html = "".join(
        f'<span class="tag">#{t}</span>' for t in tags)
    body = post.get("본문HTML", "") or ""
    title_esc = post["제목"].replace('"', '&quot;').replace('<', '&lt;').replace('>', '&gt;')
    summary_esc = post["요약"].replace('"', '&quot;').replace('<', '&lt;').replace('>', '&gt;')[:160]
    slug = post["slug"]
    date_str = post["날짜"]
    hits_url = f"https%3A%2F%2Faligomedia.co.kr%2Fblog%2F{slug}%2F"
    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>{title_esc} | 바이럴 백과사전 | 알리고미디어</title>
<meta name="description" content="{summary_esc}">
<meta name="keywords" content="{tag_meta}, 바이럴마케팅, 알리고미디어">
<meta name="robots" content="index,follow">
<link rel="canonical" href="https://aligomedia.co.kr/blog/{slug}/">
<meta property="og:type" content="article">
<meta property="og:title" content="{title_esc}">
<meta property="og:description" content="{summary_esc}">
<meta property="og:url" content="https://aligomedia.co.kr/blog/{slug}/">
<meta property="og:site_name" content="알리고미디어">
<link rel="icon" type="image/png" href="https://aligomedia.co.kr/favicon.png">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap" rel="stylesheet">
<style>
*{{box-sizing:border-box;margin:0;padding:0;}}
body{{font-family:'Noto Sans KR',sans-serif;background:#fff;color:#1a1a1a;line-height:1.75;font-size:16px;}}
header{{position:sticky;top:0;background:rgba(255,255,255,0.98);backdrop-filter:blur(10px);padding:0 5%;display:flex;justify-content:space-between;align-items:center;border-bottom:1px solid #eee;z-index:1000;height:80px;}}
.logo-area img{{height:60px;width:auto;}}
.nav-menu{{display:flex;gap:2rem;}}
.nav-menu a{{text-decoration:none;color:#1a1a1a;font-weight:700;font-size:0.9rem;padding:8px 12px;transition:0.3s;}}
.nav-menu a:hover,.nav-menu a.active{{color:#2e4a8f;}}
.post-header{{background:linear-gradient(135deg,#2e4a8f 0%,#1a2f6b 100%);color:#fff;padding:70px 5% 50px;text-align:center;}}
.post-header .tags{{display:flex;justify-content:center;flex-wrap:wrap;gap:8px;margin-bottom:20px;}}
.tag{{background:rgba(255,255,255,0.2);color:#fff;font-size:0.8rem;font-weight:700;padding:4px 12px;border-radius:20px;}}
.post-header h1{{font-size:clamp(1.6rem,3.5vw,2.5rem);font-weight:900;letter-spacing:-1px;line-height:1.3;margin-bottom:14px;max-width:800px;margin-left:auto;margin-right:auto;}}
.post-header .meta{{font-size:0.85rem;opacity:0.75;}}
.post-container{{max-width:800px;margin:0 auto;padding:60px 5% 80px;}}
.post-summary{{background:#f4f6fb;border-left:4px solid #2e4a8f;padding:18px 22px;border-radius:0 12px 12px 0;margin-bottom:40px;font-size:1rem;color:#333;font-style:italic;line-height:1.7;}}
.post-body{{font-size:1rem;line-height:1.85;color:#222;}}
.post-body h1,.post-body h2,.post-body h3{{margin:1.8em 0 0.6em;font-weight:900;letter-spacing:-0.5px;}}
.post-body p{{margin-bottom:1.2em;}}
.post-body ul,.post-body ol{{margin:1em 0 1.2em 1.5em;}}
.post-body li{{margin-bottom:0.4em;}}
.post-body blockquote{{border-left:4px solid #2e4a8f;padding:12px 18px;background:#f4f6fb;margin:1.5em 0;color:#444;}}
.post-body img{{max-width:100%;border-radius:8px;margin:1em 0;}}
.post-body strong{{color:#1a1a1a;font-weight:800;}}
.post-body .ql-align-center{{text-align:center;}}
.post-body .ql-align-right{{text-align:right;}}
.post-body .ql-align-justify{{text-align:justify;}}
.post-body p.ql-align-center{{text-align:center;}}
.post-body p.ql-align-right{{text-align:right;}}
.post-body p.ql-align-justify{{text-align:justify;}}
.post-views{{text-align:right;padding:16px 0 8px;}}
.post-footer{{border-top:1px solid #eee;padding:40px 5%;text-align:center;}}
.back-btn{{display:inline-block;background:#2e4a8f;color:#fff;font-weight:700;font-size:0.9rem;padding:12px 28px;border-radius:8px;text-decoration:none;transition:0.3s;}}
.back-btn:hover{{background:#1a2f6b;}}
footer{{padding:40px 5%;background:#fff;border-top:1px solid #eee;}}
.footer-content{{max-width:1100px;margin:0 auto;}}
.footer-info{{font-size:0.83rem;color:#555;line-height:1.8;}}
.footer-copy{{margin-top:12px;font-size:0.78rem;color:#bbb;}}
@media(max-width:768px){{.nav-menu{{gap:1rem;}}.nav-menu a{{font-size:0.8rem;padding:6px 8px;}}}}
</style>
</head>
<body>
<header>
<div class="logo-area">
<a href="https://aligomedia.co.kr"><img src="https://aligomedia.co.kr/image_4.png" alt="알리고미디어 로고" height="60"></a>
</div>
<nav class="nav-menu">
<a href="https://aligomedia.co.kr/#pr-section">언론보도대행</a>
<a href="https://aligomedia.co.kr/#review-section">리뷰마케팅</a>
<a href="https://aligomedia.co.kr/#process-section">업무프로세스</a>
<a href="/blog/" class="active">바이럴 백과사전</a>
</nav>
</header>
<section class="post-header">
<div class="tags">{tag_header_html}</div>
<h1>{post["제목"]}</h1>
<div class="meta">{date_str}</div>
</section>
<main class="post-container">
<div class="post-summary">{post["요약"]}</div>
<div class="post-body">{body}</div>
</main>
<div class="post-views">
<img src="https://hits.seeyoufarm.com/api/count/incr/badge.svg?url={hits_url}&count_bg=%232e4a8f&title_bg=%23555555&title=%EC%A1%B0%ED%9A%8C%EC%88%98&edge_flat=false" alt="조회수" height="22" loading="lazy">
</div>
<div class="post-footer">
<a href="/blog/" class="back-btn">← 목록으로 돌아가기</a>
</div>
<footer>
<div class="footer-content">
<div class="footer-info">
<div style="font-weight:800;font-size:1rem;margin-bottom:8px;">알리고미디어</div>
대표 : 박철규, 한용범 | 사업자등록번호 : 161-22-02310<br>
소재지 : 서울 마포구 양화로 64, 8층 LS-814호 | Email : helper@aligomedia.kr
<div class="footer-copy">© Aligo Media. All rights reserved.</div>
</div>
</div>
</footer>
</body>
</html>"""


def generate_blog_index_html(posts):
    """블로그 인덱스 HTML 생성 (최신순 카드 목록)"""
    if not posts:
        cards_html = """        <div class="empty-state" style="grid-column:1/-1;">
            <div class="icon">&#x270D;</div>
            <p>아직 게시글이 없습니다.<br>첫 번째 글이 곧 올라올 예정입니다!</p>
        </div>"""
    else:
        cards = []
        for p in posts:
            tags = [t.strip() for t in p["해시태그"].split(",") if t.strip()]
            tag_html = "".join(
                f'<span class="card-tag">#{t}</span>' for t in tags[:4])
            title_esc = p["제목"].replace('<', '&lt;').replace('>', '&gt;')
            summary_esc = p["요약"].replace('<', '&lt;').replace('>', '&gt;')
            # 썸네일: 본문 첫 이미지 추출
            thumb_url = extract_first_image(p.get("본문HTML", ""))
            if thumb_url:
                thumb_html = f'<div class="card-thumb"><img src="{thumb_url}" alt="{title_esc}" loading="lazy"></div>'
            else:
                thumb_html = '<div class="card-thumb-placeholder">📖</div>'
            cards.append(f"""        <article class="post-card">
            {thumb_html}
            <div class="card-body">
                <div class="card-tags">{tag_html}</div>
                <a class="card-title" href="/blog/{p["slug"]}/">{title_esc}</a>
                <p class="card-summary">{summary_esc}</p>
                <div class="card-meta">
                    <span>{p["날짜"][:10]}</span>
                    <a class="card-read" href="/blog/{p["slug"]}/">자세히 읽기 &rarr;</a>
                </div>
            </div>
        </article>""")
        cards_html = "\n".join(cards)

    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>바이럴 백과사전 | 알리고미디어</title>
<meta name="description" content="언론보도, 리뷰마케팅, 바이럴 마케팅에 관한 실전 정보를 알리고미디어가 정리합니다.">
<meta name="keywords" content="바이럴마케팅, 언론보도대행, 리뷰마케팅, 쿠팡리뷰, 보도자료, 알리고미디어">
<meta name="robots" content="index,follow">
<link rel="canonical" href="https://aligomedia.co.kr/blog/">
<meta property="og:type" content="website">
<meta property="og:title" content="바이럴 백과사전 | 알리고미디어">
<meta property="og:description" content="언론보도, 리뷰마케팅, 바이럴 마케팅에 관한 실전 정보를 알리고미디어가 정리합니다.">
<meta property="og:url" content="https://aligomedia.co.kr/blog/">
<meta property="og:site_name" content="알리고미디어">
<link rel="icon" type="image/png" href="https://aligomedia.co.kr/favicon.png">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;700;900&display=swap" rel="stylesheet">
<style>
*{{box-sizing:border-box;margin:0;padding:0;}}
html{{scroll-behavior:smooth;}}
body{{font-family:'Noto Sans KR',sans-serif;background:#fff;color:#1a1a1a;line-height:1.7;font-size:16px;}}
header{{position:sticky;top:0;background:rgba(255,255,255,0.98);backdrop-filter:blur(10px);padding:0 5%;display:flex;justify-content:space-between;align-items:center;border-bottom:1px solid #eee;z-index:1000;height:80px;}}
.logo-area img{{height:60px;width:auto;display:block;}}
.nav-menu{{display:flex;gap:2rem;}}
.nav-menu a{{text-decoration:none;color:#1a1a1a;font-weight:700;font-size:0.9rem;padding:8px 12px;transition:0.3s;}}
.nav-menu a:hover,.nav-menu a.active{{color:#2e4a8f;}}
.page-hero{{background:linear-gradient(135deg,#2e4a8f 0%,#1a2f6b 100%);color:#fff;text-align:center;padding:70px 5% 60px;}}
.page-hero h1{{font-size:clamp(1.8rem,4vw,2.8rem);font-weight:900;letter-spacing:-1px;margin-bottom:14px;}}
.page-hero p{{font-size:1.05rem;opacity:0.85;max-width:600px;margin:0 auto;}}
.blog-container{{max-width:1100px;margin:0 auto;padding:60px 5%;}}
.blog-grid{{display:grid;grid-template-columns:repeat(auto-fill,minmax(320px,1fr));gap:32px;}}
.post-card{{border:1px solid #eee;border-radius:16px;overflow:hidden;transition:box-shadow 0.3s,transform 0.3s;display:flex;flex-direction:column;}}
.post-card:hover{{box-shadow:0 8px 32px rgba(46,74,143,0.12);transform:translateY(-4px);}}
.card-thumb{{width:100%;height:200px;overflow:hidden;background:#eef1fa;}}
.card-thumb img{{width:100%;height:100%;object-fit:cover;display:block;transition:transform 0.4s;}}
.post-card:hover .card-thumb img{{transform:scale(1.05);}}
.card-thumb-placeholder{{width:100%;height:200px;background:linear-gradient(135deg,#2e4a8f22,#2e4a8f44);display:flex;align-items:center;justify-content:center;font-size:2.5rem;}}
.card-body{{padding:24px;flex:1;display:flex;flex-direction:column;}}
.card-tags{{display:flex;flex-wrap:wrap;gap:6px;margin-bottom:14px;}}
.card-tag{{background:#eef1fa;color:#2e4a8f;font-size:0.75rem;font-weight:700;padding:4px 10px;border-radius:20px;}}
.card-title{{font-size:1.05rem;font-weight:900;letter-spacing:-0.3px;margin-bottom:10px;line-height:1.4;color:#1a1a1a;text-decoration:none;display:block;}}
.card-title:hover{{color:#2e4a8f;}}
.card-summary{{font-size:0.88rem;color:#666;line-height:1.65;flex:1;margin-bottom:18px;display:-webkit-box;-webkit-line-clamp:3;-webkit-box-orient:vertical;overflow:hidden;}}
.card-meta{{font-size:0.8rem;color:#aaa;display:flex;justify-content:space-between;align-items:center;border-top:1px solid #f0f0f0;padding-top:14px;margin-top:auto;}}
.card-read{{color:#2e4a8f;font-weight:700;font-size:0.83rem;text-decoration:none;}}
.card-read:hover{{text-decoration:underline;}}
.empty-state{{text-align:center;padding:80px 20px;color:#aaa;}}
.empty-state .icon{{font-size:3.5rem;margin-bottom:16px;}}
.empty-state p{{font-size:1rem;}}
footer{{padding:50px 5%;background:#fff;border-top:1px solid #eee;}}
.footer-content{{max-width:1100px;margin:0 auto;}}
.footer-info{{font-size:0.83rem;color:#555;line-height:1.8;}}
.footer-copy{{margin-top:20px;font-size:0.78rem;color:#bbb;}}
@media(max-width:768px){{.nav-menu{{gap:1rem;}}.nav-menu a{{font-size:0.8rem;padding:6px 8px;}}.blog-grid{{grid-template-columns:1fr;}}}}
</style>
</head>
<body>
<header>
<div class="logo-area">
<a href="https://aligomedia.co.kr"><img src="https://aligomedia.co.kr/image_4.png" alt="알리고미디어 로고" height="60"></a>
</div>
<nav class="nav-menu">
<a href="https://aligomedia.co.kr/#pr-section">언론보도대행</a>
<a href="https://aligomedia.co.kr/#review-section">리뷰마케팅</a>
<a href="https://aligomedia.co.kr/#process-section">업무프로세스</a>
<a href="/blog/" class="active">바이럴 백과사전</a>
</nav>
</header>
<section class="page-hero">
<h1>&#128214; 바이럴 백과사전</h1>
<p>언론보도, 리뷰마케팅, 바이럴 마케팅의 실전 노하우를 알리고미디어가 직접 정리합니다.</p>
</section>
<main class="blog-container">
<div class="blog-grid">
{cards_html}
</div>
</main>
<footer>
<div class="footer-content">
<div class="footer-info">
<div style="font-weight:800;font-size:1rem;margin-bottom:8px;">알리고미디어</div>
대표 : 박철규, 한용범 | 사업자등록번호 : 161-22-02310<br>
소재지 : 서울 마포구 양화로 64, 8층 LS-814호 | Email : helper@aligomedia.kr
<div class="footer-copy">© Aligo Media. All rights reserved.</div>
</div>
</div>
</footer>
</body>
</html>"""


def deploy_blog_incremental(token, site_id, new_files):
    """
    Netlify 해시 기반 증분 배포.
    new_files: {path: bytes} (예: {"blog/index.html": b"...", "blog/slug/index.html": b"..."})
    기존 사이트 파일을 유지하면서 지정된 파일만 추가/갱신한다.
    """
    import hashlib as _hl
    # 유효성 검사
    if not token:
        return False, "NETLIFY_TOKEN 시크릿이 비어 있습니다. Streamlit Cloud 시크릿 설정을 확인해주세요."
    if not site_id:
        return False, "NETLIFY_SITE_ID 시크릿이 비어 있습니다. Streamlit Cloud 시크릿 설정을 확인해주세요."
    headers_auth = {"Authorization": f"Bearer {token}"}

    # 1. 최신 배포 파일 목록 가져오기
    try:
        r = requests.get(
            f"https://api.netlify.com/api/v1/sites/{site_id}/deploys?per_page=10",
            headers=headers_auth, timeout=20)
        deploys = r.json() if r.status_code == 200 else []
        latest_id = None
        for d in (deploys if isinstance(deploys, list) else []):
            if d.get("state") == "ready":
                latest_id = d["id"]
                break

        existing_files = {}
        if latest_id:
            r2 = requests.get(
                f"https://api.netlify.com/api/v1/deploys/{latest_id}/files",
                headers=headers_auth, timeout=20)
            if r2.status_code == 200:
                for f in r2.json():
                    path = f.get("id", "").lstrip("/")
                    sha = f.get("sha", "")
                    if path and sha:
                        existing_files[path] = sha
    except Exception as e:
        return False, f"기존 배포 파일 조회 실패: {e}"

    # 2. 새 파일 SHA1 계산
    new_sha = {}
    content_by_sha = {}
    for path, content in new_files.items():
        path_clean = path.lstrip("/")
        sha1 = _hl.sha1(content).hexdigest()
        new_sha[path_clean] = sha1
        content_by_sha[sha1] = (path_clean, content)

    # 3. 병합 (기존 유지 + 새 파일 덮어쓰기)
    merged = {**existing_files, **new_sha}

    # 4. 새 배포 생성
    try:
        r3 = requests.post(
            f"https://api.netlify.com/api/v1/sites/{site_id}/deploys",
            headers={**headers_auth, "Content-Type": "application/json"},
            json={"files": {f"/{k}": v for k, v in merged.items()}},
            timeout=30)
        if r3.status_code not in [200, 201]:
            return False, (f"배포 생성 실패 (HTTP {r3.status_code})\n"
                           f"site_id: {site_id}\n"
                           f"응답: {r3.text[:500]}")
        deploy_data = r3.json()
        new_deploy_id = deploy_data["id"]
        required = deploy_data.get("required", [])
    except Exception as e:
        return False, f"배포 생성 오류: {e}"

    # 5. 필요한 파일만 업로드
    upload_errors = []
    for sha1 in required:
        if sha1 not in content_by_sha:
            continue
        file_path, file_content = content_by_sha[sha1]
        try:
            ru = requests.put(
                f"https://api.netlify.com/api/v1/deploys/{new_deploy_id}/files/{file_path}",
                headers={**headers_auth, "Content-Type": "application/octet-stream"},
                data=file_content, timeout=30)
            if ru.status_code not in [200, 201]:
                upload_errors.append(f"{file_path}: {ru.status_code}")
        except Exception as e:
            upload_errors.append(f"{file_path}: {e}")

    if upload_errors:
        return False, "일부 파일 업로드 실패:\n" + "\n".join(upload_errors)
    return True, f"배포 성공 (deploy_id: {new_deploy_id})"


def upload_blog_image(token, site_id, image_bytes, original_filename):
    """
    블로그 이미지를 Netlify에 업로드하고 공개 URL을 반환.
    blog/images/{timestamp}-{filename} 경로에 저장.
    """
    ts = datetime.now().strftime("%Y%m%d%H%M%S%f")[:17]
    safe_name = re.sub(r'[^\w.\-]', '_', original_filename)
    img_path = f"blog/images/{ts}-{safe_name}"
    ok, msg = deploy_blog_incremental(token, site_id, {img_path: image_bytes})
    if ok:
        return f"https://aligomedia.co.kr/{img_path}", "성공"
    return None, msg


def process_quill_images(html_content, token, site_id):
    """
    Quill 에디터 HTML에 포함된 base64 이미지를 모두 추출해
    Netlify에 업로드하고 URL로 교체한 HTML을 반환.
    - 동일 base64는 한 번만 업로드 (중복 방지)
    - 업로드 실패한 이미지는 base64 그대로 유지
    """
    import base64 as _b64
    # base64 이미지 데이터 URI 패턴
    pattern = re.compile(r'data:image/(\w+);base64,([A-Za-z0-9+/]+=*)')
    matches = list(pattern.finditer(html_content))
    if not matches:
        return html_content, 0

    # 중복 제거 (같은 base64 → 한 번만 업로드)
    seen = {}
    for m in matches:
        key = m.group(0)
        if key not in seen:
            seen[key] = {"type": m.group(1), "b64": m.group(2), "url": None}

    # 업로드
    uploaded = 0
    for i, (data_uri, info) in enumerate(seen.items()):
        try:
            img_bytes = _b64.b64decode(info["b64"])
        except Exception:
            continue
        ext = "jpg" if info["type"] == "jpeg" else info["type"]
        fname = f"img_{i}.{ext}"
        url, _ = upload_blog_image(token, site_id, img_bytes, fname)
        if url:
            info["url"] = url
            uploaded += 1

    # 치환
    result = html_content
    for data_uri, info in seen.items():
        if info["url"]:
            result = result.replace(data_uri, info["url"])

    return result, uploaded


def blocks_to_html(blocks):
    """
    블록 목록 → 단일 HTML 문자열.
    텍스트 블록: Quill HTML 그대로 사용
    이미지 블록: <img> 태그 생성
    """
    parts = []
    for blk in blocks:
        if blk["type"] == "text":
            content = st.session_state.get(f"__vbtxt_{blk['id']}", "").strip()
            if content and content not in ("<p><br></p>", "<p></p>", ""):
                parts.append(content)
        elif blk["type"] == "image":
            url = blk.get("url", "")
            if url:
                parts.append(
                    f'<p style="text-align:center;margin:1.2em 0;">'
                    f'<img src="{url}" style="max-width:100%;border-radius:8px;" alt="이미지">'
                    f'</p>')
    return "\n".join(parts)

# ─────────────────────────────────────────────

def parse_reviews(text):
    delim = re.compile(r'^\s*(?:\((\d+)\)|(\d+)[.\)]|(\d+))\s*$', re.MULTILINE)
    markers = [(int(m.group(1) or m.group(2) or m.group(3)), m.start(), m.end()) for m in delim.finditer(text)]
    if not markers: return []
    reviews = []
    for i,(num,start,end) in enumerate(markers):
        raw = text[end:markers[i+1][1]] if i+1<len(markers) else text[end:]
        content = raw.strip()
        if content: reviews.append((num, content))
    return sorted(reviews, key=lambda x: x[0])

def create_excel(reviews):
    wb = Workbook(); ws = wb.active; ws.title = "리뷰"

    # ── 스타일 ──
    green_fill   = PatternFill("solid", start_color="70AD47", end_color="70AD47")
    h_font       = Font(bold=True, color="FFFFFF", name="맑은 고딕", size=11)
    red_font     = Font(bold=True, color="FF0000", name="맑은 고딕", size=10)
    data_font    = Font(name="맑은 고딕", size=10)
    bold_font    = Font(name="맑은 고딕", size=10, bold=True)
    center       = Alignment(horizontal="center", vertical="center", wrap_text=False)
    lw           = Alignment(horizontal="left",   vertical="top",    wrap_text=True)
    thin         = Side(style="thin", color="BFBFBF")
    border       = Border(left=thin, right=thin, top=thin, bottom=thin)

    # ── 헤더 행 (A1:E1) ──
    for col, name in [(1,"No"),(2,"별점"),(3,"리뷰내용"),(4,"사진"),(5,"상품옵션")]:
        c = ws.cell(row=1, column=col, value=name)
        c.fill = green_fill; c.font = h_font; c.alignment = center; c.border = border

    # ── F1 안내문구 (초록 배경 + 빨간 볼드) ──
    ws.merge_cells("F1:H1")
    f1 = ws["F1"]
    f1.value     = "사진은 해당 리뷰에 맞는 넘버로 동일하게 포토파일명을 변경해주세요!"
    f1.font      = red_font
    f1.alignment = Alignment(horizontal="left", vertical="center", wrap_text=False)

    # ── 컬럼 너비 ──
    ws.column_dimensions["A"].width = 6
    ws.column_dimensions["B"].width = 11
    ws.column_dimensions["C"].width = 75
    ws.column_dimensions["D"].width = 14
    ws.column_dimensions["E"].width = 14
    ws.column_dimensions["F"].width = 48

    # ── 헤더 행 높이 + AutoFilter ──
    ws.row_dimensions[1].height = 25
    ws.auto_filter.ref = "A1:E1"

    # ── 데이터 행 ──
    for i, (num, content) in enumerate(reviews, start=2):
        a  = ws.cell(row=i, column=1, value=num)
        a.font = bold_font; a.alignment = center; a.border = border

        b  = ws.cell(row=i, column=2, value="별점 5점")
        b.font = data_font; b.alignment = center; b.border = border

        cc = ws.cell(row=i, column=3, value=content)
        cc.font = data_font; cc.alignment = lw; cc.border = border

        d  = ws.cell(row=i, column=4, value=""); d.border = border
        e  = ws.cell(row=i, column=5, value=""); e.border = border

        ws.row_dimensions[i].height = max(20, min(content.count('\n') * 15 + 18, 200))

    out = io.BytesIO(); wb.save(out); out.seek(0); return out

def analyze_images_with_claude(client, image_data_list):
    if not image_data_list:
        return ""
    content = []
    for img in image_data_list:
        content.append({"type": "image", "source": {"type": "base64", "media_type": img["media_type"], "data": img["data"]}})
    content.append({"type": "text", "text": """이 제품 이미지들을 분석해서 리뷰 작성에 활용할 수 있도록 아래 항목을 상세하게 설명해줘.
1. 제품 외관 및 디자인 (색상, 형태, 크기감, 패키징)
2. 제품에 표시된 텍스트, 로고, 브랜드명, 성분 등
3. 제품의 재질감, 질감, 마감 느낌
4. 이미지에서 보이는 특징적인 요소들
5. 전반적인 제품 분위기
리뷰 작가가 실제로 제품을 써본 것처럼 묘사할 수 있도록 구체적으로 써줘. 설명만 출력해."""})
    response = client.messages.create(model="claude-sonnet-4-6", max_tokens=1000, messages=[{"role": "user", "content": content}])
    return response.content[0].text.strip()

def generate_reviews_with_claude(client, product_info, selling_points, review_count, char_count, image_data_list=None, progress_callback=None):
    import random
    BATCH_SIZE = 20
    image_description = ""
    if image_data_list:
        image_description = analyze_images_with_claude(client, image_data_list)
    persona_pool = [
        "20대 초반 여성 대학생", "20대 중반 직장 여성", "20대 후반 직장 여성",
        "20대 초반 남성 대학생", "20대 후반 남성 직장인",
        "30대 초반 주부", "30대 중반 워킹맘", "30대 후반 주부",
        "30대 초반 남성 직장인", "30대 후반 남성 직장인",
        "30대 초반 여성 자영업자", "30대 후반 여성 자영업자",
        "40대 초반 주부", "40대 중반 주부", "40대 후반 주부",
        "40대 초반 남성 회사원", "40대 중반 남성 회사원", "40대 후반 남성 회사원",
        "40대 여성 자영업자", "40대 워킹맘",
        "50대 초반 주부", "50대 후반 주부",
        "50대 초반 남성", "50대 후반 남성",
        "60대 초반 여성", "60대 후반 여성",
        "20대 여성 프리랜서", "30대 여성 교사", "40대 여성 간호사",
        "20대 남성 군인", "30대 남성 공무원", "40대 남성 자영업자",
        "30대 싱글 여성", "30대 신혼 여성", "40대 싱글 남성",
        "20대 후반 여성 간호사", "30대 여성 약사", "40대 여성 교사",
        "50대 남성 공무원", "60대 남성 은퇴자",
        "20대 남성 배달기사", "30대 남성 IT개발자", "40대 남성 의사",
        "20대 여성 헤어디자이너", "30대 여성 요가강사", "40대 여성 영양사",
        "50대 여성 교사", "60대 여성 주부", "30대 워킹대디", "40대 워킹대디"
    ]
    random.shuffle(persona_pool)
    all_personas = [persona_pool[i % len(persona_pool)] for i in range(review_count)]
    all_reviews = []
    batches = []
    for start in range(0, review_count, BATCH_SIZE):
        end = min(start + BATCH_SIZE, review_count)
        batches.append((start, end))
    for batch_idx, (start, end) in enumerate(batches):
        batch_num = batch_idx + 1
        batch_personas = all_personas[start:end]
        batch_count = end - start
        global_start_num = start + 1
        persona_text = "\n".join([f"{global_start_num + i}번 리뷰: {p}" for i, p in enumerate(batch_personas)])
        prev_summary = ""
        if all_reviews:
            recent = all_reviews[-20:]
            prev_lines = []
            for num, content in recent:
                first_line = content.split('\n')[0][:60]
                prev_lines.append(f"- {num}번: {first_line}...")
            prev_summary = "\n[이미 작성된 리뷰 도입부 (절대 유사하게 쓰지 말 것)]\n" + "\n".join(prev_lines)
        image_section = ""
        if image_description:
            image_section = f"\n[제품 이미지 분석 결과]\n{image_description}\n"
        prompt = f"""너는 실제 구매자처럼 자연스러운 한국어 리뷰를 쓰는 전문 작가야.
[제품 정보]
{product_info}
[소구점]
{selling_points if selling_points else "없음"}
{image_section}
[작성 조건]
- 이번 배치: {global_start_num}번 ~ {end}번 리뷰 (총 {batch_count}개)
- 리뷰당 글자 수: 약 {char_count}자 내외
[페르소나 배정]
{persona_text}
{prev_summary}
[필수 규칙]
1. 번호는 {global_start_num}부터 시작, 각 리뷰는 "번호." 한 줄 후 리뷰 내용
2. 그림 이모지 절대 사용 금지
3. 텍스트 감성 표현 자연스럽게 허용 (ㅎㅎ, ㅋㅋ 등)
4. 리뷰 간 표현·문장구조 절대 중복 금지
5. 페르소나에 맞는 실제 사람 말투 사용
정확히 {batch_count}개 리뷰를 작성해줘. 설명이나 부연 없이 리뷰만 출력해."""
        response = client.messages.create(model="claude-sonnet-4-6", max_tokens=8000, messages=[{"role": "user", "content": prompt}])
        raw_text = response.content[0].text.strip()
        batch_reviews = parse_generated_reviews(raw_text)
        if batch_reviews:
            all_reviews.extend(batch_reviews)
        if progress_callback:
            progress_callback(batch_idx + 1, len(batches), len(all_reviews))
    return all_reviews

def parse_generated_reviews(text):
    lines = text.split('\n')
    reviews = []
    current_num = None
    current_lines = []
    for line in lines:
        stripped = line.strip()
        num_match = re.match(r'^(\d+)[.\)]?\s*$', stripped)
        if num_match:
            if current_num is not None and current_lines:
                content = '\n'.join(current_lines).strip()
                if content:
                    reviews.append((current_num, content))
            current_num = int(num_match.group(1))
            current_lines = []
        else:
            if current_num is not None:
                current_lines.append(line)
    if current_num is not None and current_lines:
        content = '\n'.join(current_lines).strip()
        if content:
            reviews.append((current_num, content))
    return sorted(reviews, key=lambda x: x[0])

def parse_match_response(text):
    mc, mb = "미등록", "미등록"
    for line in text.split('\n'):
        line = line.strip()
        if not line or ':' not in line:
            continue
        # 상품코드 파싱 — "상품코드" 포함 줄에서 콜론 뒤 값 추출
        if '상품코드' in line:
            val = line.split(':', 1)[1]
            val = re.sub(r'[*_`\[\]()\s]', '', val).strip()
            if val and val not in ('미등록', '없음', ''):
                mc = val
            elif val in ('미등록', '없음'):
                mc = '미등록'
        # 브랜드 파싱 — elif로 같은 줄 중복 처리 방지
        elif '브랜드' in line:
            val = line.split(':', 1)[1]
            val = re.sub(r'[*_`\[\]]', '', val).strip()
            if val and val not in ('미등록', '없음', ''):
                mb = val
            elif val in ('미등록', '없음'):
                mb = '미등록'
    return mc, mb


# ==================== 함소아 보고서 관련 ====================

HAMSOA_SHEET_URL = "https://docs.google.com/spreadsheets/d/1yozxvC3iXhCkbC3yXf5ad6PHaaZC3yQEvEXhpRcnGwc/"

def get_hamsoa_sheet(worksheet_name):
    try:
        scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
        try:
            creds_dict = json.loads(st.secrets["GOOGLE_CREDENTIALS"])
            creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
        except:
            BASE_DIR = os.path.dirname(os.path.abspath(__file__))
            creds = ServiceAccountCredentials.from_json_keyfile_name(
                os.path.join(BASE_DIR, 'service_account.json'), scope)
        client_gs = gspread.authorize(creds)
        return client_gs.open_by_url(HAMSOA_SHEET_URL).worksheet(worksheet_name)
    except Exception as e:
        st.error(f"함소아 시트 연결 실패 ({worksheet_name}): {e}")
        return None

def find_current_block_start(rows, no_col=0):
    """A열 번호가 1로 재시작되는 마지막 위치 반환"""
    last_start = 0
    prev_num = 0
    for i, row in enumerate(rows):
        if len(row) > no_col:
            cell = str(row[no_col]).strip()
            try:
                num = int(cell)
                if num == 1 and prev_num >= 5:
                    last_start = i
                prev_num = num
            except:
                pass
    return last_start

def _get_field(rec, keys):
    for k in keys:
        if k in rec and str(rec[k]).strip():
            return str(rec[k]).strip()
    return ''

def parse_competitor_sheet(sheet):
    all_values = sheet.get_all_values()
    if not all_values:
        return [], {}

    # 헤더 행 찾기
    header_row_idx = 0
    for i, row in enumerate(all_values):
        joined = ' '.join(str(c) for c in row)
        if ('경쟁사' in joined or 'NO' in joined.upper()) and '매체사' in joined and '발행' in joined:
            header_row_idx = i
            break

    headers = [str(h).strip() for h in all_values[header_row_idx]]
    data_rows = all_values[header_row_idx + 1:]

    block_start = find_current_block_start(data_rows, 0)
    current_rows = data_rows[block_start:]

    records = []
    for row in current_rows:
        if len(row) > 1 and str(row[0]).strip().isdigit():
            record = {}
            for j, h in enumerate(headers):
                record[h] = row[j].strip() if j < len(row) else ''
            records.append(record)

    # 병원별 통계
    hospital_counts = {}
    strategy_by_hospital = {}
    for rec in records:
        comp = _get_field(rec, ['경쟁사', '병원', '업체'])
        if not comp:
            continue
        hospital_counts[comp] = hospital_counts.get(comp, 0) + 1
        strategy = _get_field(rec, ['전략유형', '전략 유형', '전략'])
        if comp not in strategy_by_hospital:
            strategy_by_hospital[comp] = {}
        if strategy:
            strategy_by_hospital[comp][strategy] = strategy_by_hospital[comp].get(strategy, 0) + 1

    return records, {
        'hospital_counts': hospital_counts,
        'strategy_by_hospital': strategy_by_hospital,
        'headers': headers,
    }

def parse_hamsoa_sheet(sheet):
    all_values = sheet.get_all_values()
    if not all_values:
        return [], []

    articles = []
    billing = []

    article_header_idx = -1
    billing_header_idx = -1

    for i, row in enumerate(all_values):
        joined = ' '.join(str(c) for c in row)
        if '발행일' in joined and '구분' in joined and '제목' in joined and article_header_idx < 0:
            article_header_idx = i
        if '매체사' in joined and '분류' in joined and ('건당' in joined or '견적' in joined) and billing_header_idx < 0:
            billing_header_idx = i

    if article_header_idx >= 0:
        art_headers = [str(h).strip() for h in all_values[article_header_idx]]
        end_idx = billing_header_idx if 0 < billing_header_idx > article_header_idx else len(all_values)
        for i, row in enumerate(all_values[article_header_idx + 1:end_idx], start=article_header_idx + 1):
            if len(row) > 1 and str(row[0]).strip() and str(row[0]).strip() not in ['발행일', '']:
                record = {}
                for j, h in enumerate(art_headers):
                    record[h] = row[j].strip() if j < len(row) else ''
                articles.append(record)

    if billing_header_idx >= 0:
        bill_headers = [str(h).strip() for h in all_values[billing_header_idx]]
        for row in all_values[billing_header_idx + 1:]:
            if len(row) > 1 and str(row[0]).strip():
                cell0 = str(row[0]).strip().upper()
                if 'TOTAL' in cell0 or 'TOAL' in cell0:
                    continue
                record = {}
                has_data = False
                for j, h in enumerate(bill_headers):
                    val = row[j].strip() if j < len(row) else ''
                    record[h] = val
                    if val and j > 0:
                        has_data = True
                if has_data:
                    billing.append(record)

    return articles, billing

def make_bar_chart(hospital_counts, hamsoa_count):
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    try:
        plt.rcParams['font.family'] = 'Malgun Gothic'
    except:
        try:
            plt.rcParams['font.family'] = 'NanumGothic'
        except:
            pass
    plt.rcParams['axes.unicode_minus'] = False

    color_map = {
        '함소아 한의원': '#4472C4',
        '자생한방병원': '#FF4444',
        '폴리한의원': '#FFC000',
        '아이누리한의원': '#70AD47',
        '꽃피는 한의원': '#ED7D31',
        '해아림한의원': '#4BACC6',
        '헤아림한의원': '#4BACC6',
    }

    all_data = {'함소아 한의원': hamsoa_count}
    all_data.update(hospital_counts)

    sorted_names = ['함소아 한의원'] + sorted(
        [h for h in hospital_counts], key=lambda x: hospital_counts.get(x, 0), reverse=True)
    counts = [all_data.get(h, 0) for h in sorted_names]
    colors = [color_map.get(h, '#888888') for h in sorted_names]

    fig, ax = plt.subplots(figsize=(11, 5))
    bars = ax.bar(range(len(sorted_names)), counts, color=colors, width=0.5, edgecolor='white')

    ax.set_xticks(range(len(sorted_names)))
    ax.set_xticklabels(sorted_names, fontsize=10)
    ax.set_title("각 병원별 보도자료 배포 수량", fontsize=13, fontweight='bold', pad=15)
    ax.yaxis.grid(True, linestyle='--', alpha=0.5)
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    # 범례
    legend_patches = [plt.Rectangle((0, 0), 1, 1, color=color_map.get(n, '#888888')) for n in sorted_names]
    ax.legend(legend_patches, sorted_names, loc='upper right', fontsize=8, ncol=2)

    for bar, count in zip(bars, counts):
        if count > 0:
            ax.text(bar.get_x() + bar.get_width() / 2., bar.get_height() + 0.3,
                    str(count), ha='center', va='bottom', fontsize=10, fontweight='bold')

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    plt.close()
    return buf

def make_strategy_chart(strategy_by_hospital):
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt

    try:
        plt.rcParams['font.family'] = 'Malgun Gothic'
    except:
        try:
            plt.rcParams['font.family'] = 'NanumGothic'
        except:
            pass
    plt.rcParams['axes.unicode_minus'] = False

    if not strategy_by_hospital:
        fig, ax = plt.subplots(figsize=(10, 4))
        ax.text(0.5, 0.5, '데이터 없음', ha='center', va='center', fontsize=14)
        ax.axis('off')
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        plt.close()
        return buf

    all_strategies = set()
    for strats in strategy_by_hospital.values():
        all_strategies.update(strats.keys())
    strategy_list = sorted(all_strategies)

    hospitals = list(strategy_by_hospital.keys())
    x = list(range(len(hospitals)))
    width = 0.7 / max(len(strategy_list), 1)

    strat_colors = {
        '브랜드 강화형': '#4472C4',
        '질환 타깃형': '#FF4444',
        '시술 중심형': '#FFC000',
        '마케팅 형': '#70AD47',
        '마케팅형': '#70AD47',
    }

    fig, ax = plt.subplots(figsize=(11, 5))
    for i, strategy in enumerate(strategy_list):
        counts = [strategy_by_hospital[h].get(strategy, 0) for h in hospitals]
        offset = (i - len(strategy_list) / 2 + 0.5) * width
        color = strat_colors.get(strategy, f'C{i}')
        ax.bar([xi + offset for xi in x], counts, width=width * 0.9, label=strategy, color=color)

    ax.set_xticks(x)
    ax.set_xticklabels(hospitals, fontsize=9)
    ax.set_title("경쟁사 전략유형 별 집계 현황", fontsize=13, fontweight='bold', pad=15)
    ax.legend(loc='upper right', fontsize=9)
    ax.yaxis.grid(True, linestyle='--', alpha=0.5)
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    plt.tight_layout()
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    plt.close()
    return buf

def _set_cell_bg(cell, hex_color):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    for child in list(tcPr):
        tag = child.tag
        if 'Fill' in tag or 'fill' in tag:
            tcPr.remove(child)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color.replace('#', '').upper())
    tcPr.insert(0, solidFill)

def _style_cell(cell, text, bg_hex=None, fg_hex=None, font_size=8,
                bold=False, center=False, italic=False):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT

    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if fg_hex:
        h = fg_hex.replace('#', '')
        run.font.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    if bg_hex:
        _set_cell_bg(cell, bg_hex)

def parse_jasaeng_strategy(raw_text):
    """자생한방병원 전략 데이터 파싱
    형식:
        질환 타겟형
        비염 1
        일자목 1
        시술 중심형
        ...
    반환: strategies dict, order list, totals dict, grand_total int
    """
    strategies = {}
    strategy_order = []
    current_strategy = None

    for line in raw_text.strip().split('\n'):
        line = line.strip()
        if not line:
            continue
        # 줄 끝 숫자 추출 (공백 유무 상관없이: "잔기침1" or "비염 1" 모두 처리)
        m = re.match(r'^(.+?)\s*(\d+)\s*$', line)
        if m:
            keyword, count = m.group(1).strip(), int(m.group(2))
            if current_strategy is not None:
                strategies[current_strategy].append((keyword, count))
            continue
        # 전략유형 헤더 (숫자로 끝나지 않는 줄)
        current_strategy = line
        if current_strategy not in strategies:
            strategies[current_strategy] = []
            strategy_order.append(current_strategy)

    totals = {s: sum(c for _, c in items) for s, items in strategies.items()}
    grand_total = sum(totals.values())
    return strategies, strategy_order, totals, grand_total


def generate_quantity_analysis(hospital_counts, hamsoa_count, jasaeng_total, report_date):
    """병원별 발행 수량 비교 분석 텍스트 생성 (이미지 1 형식)
    반환: list of str (각 불릿 텍스트)
    """
    others = {h: c for h, c in hospital_counts.items()
              if '자생' not in h and '함소아' not in h}
    combined = list(others.items())
    combined.append(('함소아한의원', hamsoa_count))
    combined = sorted(combined, key=lambda x: x[1], reverse=True)

    positive = [(h, c) for h, c in combined if c > 0]
    zero_hosps = [h for h, c in combined if c == 0]

    texts = []

    # 불릿 1: 자생 1위
    texts.append(
        f"{report_date} 기준, 자생한방병원은 총 {jasaeng_total}건의 보도자료를 언론에 배포하며, "
        f"전체 병원 중 가장 많은 기사 발행 수를 기록했습니다."
    )

    # 불릿 2: 나머지 병원들 순서대로
    if positive:
        first_h, first_c = positive[0]
        parts2 = [f"{first_h}은 총 {first_c}건을 발행하며 그 뒤를 이었고"]
        for h, c in positive[1:]:
            parts2.append(f"{h}이 총 {c}건 발행 하였으며")
        t2 = ", ".join(parts2)
        if zero_hosps:
            t2 += f", {', '.join(zero_hosps)}은 발행 수량이 0건으로 나타났습니다."
        else:
            t2 += "."
        texts.append(t2)

    # 불릿 3: 자생 vs 함소아 비율
    if jasaeng_total > 0 and hamsoa_count > 0:
        ratio = round(hamsoa_count / jasaeng_total * 100, 1)
        times = round(jasaeng_total / hamsoa_count)
        texts.append(
            f"자생한방병원은 총 {jasaeng_total}건, 함소아한의원은 {hamsoa_count}건의 언론보도를 발행하였으며, "
            f"함소아한의원의 발행 수는 자생한방병원의 약 {ratio}% 수준으로, "
            f"양 기관 간 보도자료 운영 규모에서 {times}배 이상의 격차를 보였습니다."
        )

    return texts


def generate_strategy_analysis(jasaeng_strategies, jasaeng_order, jasaeng_totals, jasaeng_grand_total,
                                competitor_records, hospital_counts):
    """전략유형별 세부 분석 텍스트 생성 (이미지 2 형식)
    반환: list of str (각 불릿 텍스트)
    """
    texts = []

    # 자생한방병원 전략 breakdown
    strat_parts = []
    for s in jasaeng_order:
        t = jasaeng_totals.get(s, 0)
        strat_parts.append(f"'{s}'이 {t}건")
    jasaeng_strat_str = ", ".join(strat_parts)

    texts.append(
        f"총 {jasaeng_grand_total}건의 자생한방병원 기사 중, {jasaeng_strat_str}으로 확인되며, "
        f"질환 관련 정보성 기사와 브랜드 이미지 제고를 중심으로 전략적 콘텐츠 배포를 함께 다뤄진 것으로 분석됩니다. "
        f"특히 이슈가 발생한 키워드를 중점으로 칼럼 기사를 배포하여 최대한 많은 사람들에게 기사를 "
        f"노출시킬 수 있도록 집중한 것으로 확인됩니다."
    )

    # 경쟁사 스프레드시트 병원별 전략유형 집계
    hosp_strategies = {}
    for rec in competitor_records:
        h = _get_field(rec, ['병원', '경쟁사', 'hospital', '병 원'])
        s = _get_field(rec, ['전략유형', '전략 유형', 'strategy', '유형'])
        if h and s:
            hosp_strategies.setdefault(h, {})
            hosp_strategies[h][s] = hosp_strategies[h].get(s, 0) + 1

    sorted_hosps = sorted(hospital_counts.items(), key=lambda x: x[1], reverse=True)
    active = [(h, c) for h, c in sorted_hosps if c > 0 and '자생' not in h]
    inactive = [h for h, c in sorted_hosps if c == 0 and '자생' not in h]

    for h, c in active:
        strat = hosp_strategies.get(h, {})
        if strat:
            sorted_strat = sorted(strat.items(), key=lambda x: x[1], reverse=True)
            strat_desc = ", ".join(f"{s}이 {n}건" for s, n in sorted_strat)
            main_s = sorted_strat[0][0]
            texts.append(
                f"{h}은 총 {c}건의 기사 중 {strat_desc}으로 구성되어, "
                f"'{main_s}' 전략을 중심으로 운영한 것으로 분석됩니다."
            )
        else:
            texts.append(f"{h}은 총 {c}건의 기사를 발행하였습니다.")

    if inactive:
        inactive_str = ", ".join(inactive)
        texts.append(
            f"반면, {inactive_str}은 집계된 기사가 없는 상태로, 언론 노출 자체가 미비합니다."
        )

    return texts


def generate_ai_report_text(client, jasaeng_totals, jasaeng_order, jasaeng_grand_total,
                             jasaeng_strategies, hospital_counts, hamsoa_count,
                             competitor_records, report_date):
    """Claude API로 전문 분석가 어조의 보고서 텍스트 생성"""

    # 자생 전략 요약
    strat_lines = []
    for s in jasaeng_order:
        t = jasaeng_totals.get(s, 0)
        items = jasaeng_strategies.get(s, [])
        top_keywords = ", ".join(f"{kw}({cnt}건)" for kw, cnt in sorted(items, key=lambda x: x[1], reverse=True)[:5])
        strat_lines.append(f"  - {s}: {t}건" + (f" (주요: {top_keywords})" if top_keywords else ""))

    # 경쟁사 전략유형 집계
    hosp_strategies = {}
    for rec in competitor_records:
        h = _get_field(rec, ['병원', '경쟝사', 'hospital', '병 원'])
        s = _get_field(rec, ['전략유형', '전략 유형', 'strategy', '유형'])
        if h and s:
            hosp_strategies.setdefault(h, {})
            hosp_strategies[h][s] = hosp_strategies[h].get(s, 0) + 1

    hosp_lines = []
    for h, c in sorted(hospital_counts.items(), key=lambda x: x[1], reverse=True):
        strat = hosp_strategies.get(h, {})
        strat_str = ", ".join(f"{s} {n}건" for s, n in sorted(strat.items(), key=lambda x: x[1], reverse=True)) if strat else "전략유형 미분류"
        hosp_lines.append(f"  - {h}: {c}건 ({strat_str})")
    hosp_lines.append(f"  - 함소아한의원: {hamsoa_count}건")

    prompt = f"""당신은 의료 PR·미디어 전략 분석 전문가입니다.
아래 데이터를 바탕으로 경쟁사 언론보도 동향 분석 보고서 텍스트를 작성해주세요.

[분석 기준일]
{report_date}

[자생한방병원 전략유형별 기사 현황] (총 {jasaeng_grand_total}건)
{chr(10).join(strat_lines)}

[경쟁사 및 함소아한의원 기사 발행 현황]
{chr(10).join(hosp_lines)}

[작성 요청]
아래 두 섹션을 각각 불릿 3개씩 작성해주세요.

**섹션1: 병원별 언론보도 발행 수량 비교 분석**
- 자생한방병원이 압도적 1위임을 전문적으로 표현
- 다른 병원들과의 순위/수량 비교 (함소아 포함)
- 자생 대비 함소아 발행 비율·격차를 수치로 분석

**섹션2: 자생한방병원 전략유형별 세부 분석 + 경쟝사 비교**
- 자생한방병원 전략유형 구성 비율과 특징 분석 (어떤 전략이 주력인지, 왜 그런지)
- 해아림·아이누리 등 활성 경쟝사 전략 특징 분석
- 발행 저조 병원에 대한 시사점

[작성 조건]
- 전문 분석가 어조, 객관적·날카로운 시각
- 각 불릿은 2~3문장 이내
- 불릿은 반드시 "➡ " 로 시작
- 섹션 구분은 "=== 섹션1 ===" "=== 섹션2 ===" 형식으로 표시
- 한국어로 작성"""

    message = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=2000,
        messages=[{"role": "user", "content": prompt}]
    )
    return message.content[0].text


def generate_hamsoa_ppt(competitor_records, meta, hamsoa_articles, billing_data,
                         report_date, report_month, hamsoa_article_count):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    NAVY_RGB = RGBColor(0x1A, 0x2B, 0x4A)
    WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_RGB = RGBColor(0xAA, 0xAA, 0xAA)
    NAVY_HEX = "1A2B4A"
    BEIGE_HEX = "F5E3BA"
    WHITE_HEX = "FFFFFF"
    ALT_HEX = "F2F2F2"

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def new_slide():
        return prs.slides.add_slide(blank)

    def set_bg(slide, hex_color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        h = hex_color.replace('#', '')
        fill.fore_color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def txt(slide, text, left, top, width, height,
            size=11, bold=False, rgb=None, align=PP_ALIGN.LEFT, italic=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if rgb:
            run.font.color.rgb = rgb
        return box

    def aligo_mark(slide):
        txt(slide, "ALIGO", 12.1, 0.1, 1.1, 0.28, size=8, rgb=GRAY_RGB, align=PP_ALIGN.RIGHT)
        txt(slide, "MEDIA", 12.1, 0.35, 1.1, 0.28, size=8, rgb=GRAY_RGB, align=PP_ALIGN.RIGHT)

    def section_title(slide, title_text):
        txt(slide, title_text, 0.5, 0.35, 12.5, 0.55, size=14, bold=True, rgb=NAVY_RGB)
        line = slide.shapes.add_connector(1, Inches(0.5), Inches(0.95), Inches(12.8), Inches(0.95))
        line.line.color.rgb = NAVY_RGB
        line.line.width = Pt(1.2)

    def bullet_box(slide, bullets, left, top, width, height):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
        box.line.color.rgb = RGBColor(0xC5, 0xD8, 0xF0)
        box.line.width = Pt(0.75)
        tf = box.text_frame
        tf.word_wrap = True
        for i, bullet_text in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.space_before = Pt(5)
            run = para.add_run()
            run.text = f"➡  {bullet_text}"
            run.font.size = Pt(9)
            run.font.color.rgb = NAVY_RGB

    def add_pic(slide, img_buf, left, top, width, height):
        img_buf.seek(0)
        slide.shapes.add_picture(img_buf, Inches(left), Inches(top), Inches(width), Inches(height))

    def add_table(slide, data, col_widths, left, top, max_height=5.8):
        rows_n = len(data)
        cols_n = len(data[0]) if data else 1
        row_h = min(max_height / rows_n, 0.55)
        tbl_h = row_h * rows_n
        tbl = slide.shapes.add_table(
            rows_n, cols_n,
            Inches(left), Inches(top),
            Inches(sum(col_widths)), Inches(tbl_h)
        ).table
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
        for ri, row_data in enumerate(data):
            for ci, val in enumerate(row_data):
                cell = tbl.cell(ri, ci)
                is_header = ri == 0
                is_total = ri == rows_n - 1 and str(data[ri][0]).upper().startswith('TOTAL')
                is_alt = ri % 2 == 0 and not is_header

                if is_header:
                    bg, fg = BEIGE_HEX, NAVY_HEX
                elif is_total:
                    bg, fg = NAVY_HEX, WHITE_HEX
                elif is_alt:
                    bg, fg = ALT_HEX, "333333"
                else:
                    bg, fg = WHITE_HEX, "333333"

                # 예정 행 빨간 글씨
                if not is_header and '예정' in str(val):
                    fg = "CC0000"

                center_cols = {0, 1, 2, 3, 5, 6}  # 보통 날짜/숫자 컬럼
                _style_cell(cell, val, bg_hex=bg, fg_hex=fg,
                            font_size=9 if is_header else 8,
                            bold=is_header or is_total,
                            center=(ci in center_cols))
        return tbl

    # ── SLIDE 1: Cover ──────────────────────────────────────
    s1 = new_slide()
    set_bg(s1, NAVY_HEX)
    txt(s1, "ALIGO", 11.9, 0.18, 1.2, 0.3, size=9, rgb=GRAY_RGB, align=PP_ALIGN.RIGHT)
    txt(s1, "MEDIA", 11.9, 0.46, 1.2, 0.3, size=9, rgb=GRAY_RGB, align=PP_ALIGN.RIGHT)
    txt(s1, "함소아 한의원", 0.8, 1.8, 11.7, 2.0,
        size=58, bold=True, rgb=WHITE_RGB, align=PP_ALIGN.CENTER)
    txt(s1, "REPORT", 0.5, 4.9, 8.0, 1.6, size=62, bold=True, rgb=WHITE_RGB)
    txt(s1, report_date, 7.0, 5.8, 6.0, 0.55, size=14, bold=True, rgb=WHITE_RGB, align=PP_ALIGN.RIGHT)
    txt(s1, "종합 리포트", 7.0, 6.32, 6.0, 0.4, size=11, rgb=WHITE_RGB, align=PP_ALIGN.RIGHT)

    # ── SLIDE 2: Contents ────────────────────────────────────
    s2 = new_slide()
    set_bg(s2, "F0F0F0")
    aligo_mark(s2)
    txt(s2, f"{report_month} 종합 리포트", 0.7, 0.32, 8.0, 0.42, size=12, bold=True, rgb=NAVY_RGB)
    txt(s2, "CONTENTS", 0.7, 0.7, 11.5, 1.5, size=72, bold=True, rgb=NAVY_RGB)

    contents_items = [
        ("01", "언론보도 발행 수량 비교 (함소아 vs 경쟁사)"),
        ("02", "경쟁사 기사 발행 정보 상세 분석"),
        ("03", "병원별 기사 발행 현황 요약"),
        ("04", "함소아 기사 발행 현황 요약"),
        ("05", "브릿지경제 지면보도 현황"),
        ("06", "함소아 기사 집행 금액 정산표"),
    ]
    for i, (num, label) in enumerate(contents_items):
        col_idx = i // 3
        row_idx = i % 3
        bx = 0.7 + col_idx * 6.4
        by = 2.55 + row_idx * 1.25

        nb = s2.shapes.add_shape(1, Inches(bx), Inches(by), Inches(0.52), Inches(0.42))
        nb.fill.solid()
        nb.fill.fore_color.rgb = NAVY_RGB
        nb.line.fill.background()
        np_ = nb.text_frame.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = num
        nr.font.size = Pt(12)
        nr.font.bold = True
        nr.font.color.rgb = WHITE_RGB

        txt(s2, label, bx + 0.62, by + 0.02, 5.6, 0.42, size=11, rgb=NAVY_RGB)

    # ── SLIDE 3: 언론보도 발행 수량 비교 ─────────────────────
    s3 = new_slide()
    aligo_mark(s3)
    section_title(s3, "1. 언론보도 발행 수량 비교 (함소아 vs 경쟁사)")

    hospital_counts = meta.get('hospital_counts', {})
    chart_buf = make_bar_chart(hospital_counts, hamsoa_article_count)
    add_pic(s3, chart_buf, 0.4, 1.05, 9.3, 4.6)

    total_comp = sum(hospital_counts.values())
    top_h = max(hospital_counts, key=hospital_counts.get) if hospital_counts else ''
    top_c = hospital_counts.get(top_h, 0)

    bullets_s3 = []
    if top_h:
        bullets_s3.append(f"{report_date} 기준, {top_h}은(는) 총 {top_c}건으로 전체 병원 중 가장 많은 기사 발행 수를 기록했습니다.")
    bullets_s3.append(f"함소아한의원은 총 {hamsoa_article_count}건의 언론보도를 발행하였습니다.")
    for h, c in sorted(hospital_counts.items(), key=lambda x: x[1], reverse=True)[:2]:
        if h != top_h:
            bullets_s3.append(f"{h}: {c}건 발행")
    bullet_box(s3, bullets_s3[:3], 9.85, 1.05, 3.2, 4.6)

    # ── SLIDES 4+: 경쟁사 기사 발행 정보 상세 분석 ──────────
    COMP_KEY_MAP = {
        '경쟁사': ['경쟁사', '병원'],
        '매체사': ['매체사', '매체'],
        '발행일': ['발행일자', '발행일', '날짜'],
        '주요키워드': ['메인키워드', '주요키워드', '키워드', '메인 키워드'],
        '제목': ['제목', '기사제목'],
        '월간 검색량': ['월간 검색량', '월간검색량', '검색량'],
    }
    comp_header_row = ['NO.', '경쟁사', '매체사', '발행일', '주요키워드', '제목', '월간 검색량']
    comp_col_w = [0.45, 1.2, 1.2, 0.95, 1.0, 5.6, 1.0]

    ROWS_PER = 10
    total_rec = len(competitor_records)
    n_comp_slides = max(1, -(-total_rec // ROWS_PER))

    for si in range(n_comp_slides):
        sld = new_slide()
        aligo_mark(sld)
        section_title(sld, f"2. 경쟁사 기사 발행 정보 상세 분석 ({si + 1})")

        batch = competitor_records[si * ROWS_PER: (si + 1) * ROWS_PER]
        table_data = [comp_header_row]
        for k, rec in enumerate(batch):
            table_data.append([
                str(si * ROWS_PER + k + 1),
                _get_field(rec, COMP_KEY_MAP['경쟁사']),
                _get_field(rec, COMP_KEY_MAP['매체사']),
                _get_field(rec, COMP_KEY_MAP['발행일']),
                _get_field(rec, COMP_KEY_MAP['주요키워드']),
                _get_field(rec, COMP_KEY_MAP['제목']),
                _get_field(rec, COMP_KEY_MAP['월간 검색량']),
            ])
        add_table(sld, table_data, comp_col_w, left=0.35, top=1.1)

    # 경쟁사 분석 텍스트 슬라이드
    s_comp_txt = new_slide()
    aligo_mark(s_comp_txt)
    section_title(s_comp_txt, f"2. 경쟁사 기사 발행 정보 상세 분석 ({n_comp_slides + 1})")

    top_hospitals = sorted(hospital_counts.items(), key=lambda x: x[1], reverse=True)
    analysis_bullets = [f"총 {total_comp}건의 경쟁사 기사 중, {top_h}이(가) 총 {top_c}건으로 가장 높은 기사 발행 빈도를 보였습니다." if top_h else "데이터를 분석하고 있습니다."]
    for h, c in top_hospitals[:4]:
        if c > 0:
            strats = meta.get('strategy_by_hospital', {}).get(h, {})
            top_s = max(strats, key=strats.get) if strats else ''
            analysis_bullets.append(f"{h}: {c}건 ({top_s} 중심)" if top_s else f"{h}: {c}건 발행")
    bullet_box(s_comp_txt, analysis_bullets[:5], 0.5, 1.2, 12.3, 3.5)

    # ── 병원별 기사 발행 현황 ────────────────────────────────
    s_hosp = new_slide()
    aligo_mark(s_hosp)
    section_title(s_hosp, "3. 병원별 기사 발행 현황 요약")

    strat_chart_buf = make_strategy_chart(meta.get('strategy_by_hospital', {}))
    add_pic(s_hosp, strat_chart_buf, 0.4, 1.05, 9.5, 5.0)

    hosp_bullets = []
    for h, strats in sorted(meta.get('strategy_by_hospital', {}).items(),
                             key=lambda x: sum(x[1].values()), reverse=True)[:4]:
        total = sum(strats.values())
        top_s = max(strats, key=strats.get) if strats else ''
        hosp_bullets.append(f"{h}: 총 {total}건 ({top_s} {strats.get(top_s, 0)}건)")
    if hosp_bullets:
        bullet_box(s_hosp, hosp_bullets, 10.05, 1.05, 3.1, 5.0)

    # ── 함소아 기사 발행 현황 ────────────────────────────────
    ART_KEYS = ['발행일', '구분', '제목', '매체사', '메인 키워드', '검색량', '진행 현황', '본문 요약']
    ART_KEY_ALT = ['발행일', '구분', '제목', '매체사', '메인키워드', '검색량', '진행현황', '본문요약']
    ART_HDR = ['발행일', '구분', '제목', '매체사', '메인 키워드', '검색량', '진행현황', '본문 요약']
    ART_COL_W = [0.95, 0.75, 2.5, 1.05, 1.05, 0.65, 0.85, 4.7]

    ART_ROWS = 7
    total_arts = len(hamsoa_articles)
    n_art_slides = max(1, -(-total_arts // ART_ROWS))

    for si in range(n_art_slides):
        sld = new_slide()
        aligo_mark(sld)
        section_title(sld, f"4. 함소아 기사 발행 현황 요약 ({si + 1})")

        batch = hamsoa_articles[si * ART_ROWS: (si + 1) * ART_ROWS]
        table_data = [ART_HDR]
        for art in batch:
            row = []
            for k, k_alt in zip(ART_KEYS, ART_KEY_ALT):
                val = art.get(k, art.get(k_alt, art.get(k.replace(' ', ''), '')))
                row.append(val)
            table_data.append(row)
        add_table(sld, table_data, ART_COL_W, left=0.35, top=1.1)

    # 함소아 기사 분석 텍스트
    s_art_txt = new_slide()
    aligo_mark(s_art_txt)
    section_title(s_art_txt, f"4. 함소아 기사 발행 현황 요약 ({n_art_slides + 1})")

    completed = [a for a in hamsoa_articles
                 if '완료' in str(a.get('진행 현황', a.get('진행현황', '')))]
    planned = [a for a in hamsoa_articles
               if '예정' in str(a.get('진행 현황', a.get('진행현황', '')))]

    art_bullets = [
        f"{report_date} 기준, 총 {total_arts}건의 기사 중 {len(completed)}건은 발행 완료입니다.(기획기사 및 명의칼럼 포함)",
    ]
    if planned:
        art_bullets.append(f"예정된 기사: {len(planned)}건")

    keyword_counter = {}
    for a in hamsoa_articles:
        kw = a.get('메인 키워드', a.get('메인키워드', ''))
        if kw:
            keyword_counter[kw] = keyword_counter.get(kw, 0) + 1

    bullet_box(s_art_txt, art_bullets, 0.5, 1.2, 12.3, 3.5)

    # ── 브릿지경제 지면보도 ──────────────────────────────────
    s_bridge = new_slide()
    aligo_mark(s_bridge)
    section_title(s_bridge, "5. 함소아 브릿지경제 지면보도 현황")

    bridge_arts = [a for a in hamsoa_articles
                   if '브릿지' in str(a.get('매체사', ''))]

    bridge_data = [['NO.', '게재일', '키워드', '기사 제목']]
    for i, art in enumerate(bridge_arts):
        bridge_data.append([
            str(i + 1),
            art.get('발행일', ''),
            art.get('메인 키워드', art.get('메인키워드', '')),
            art.get('제목', ''),
        ])

    if len(bridge_data) > 1:
        add_table(s_bridge, bridge_data, [0.5, 1.3, 1.5, 9.2], left=0.35, top=1.1, max_height=4.0)
    else:
        txt(s_bridge, "브릿지경제 지면보도 데이터가 없습니다. (매체사 열에 '브릿지경제'로 입력된 항목이 있으면 자동 집계됩니다.)",
            0.5, 2.0, 12.0, 1.0, size=11, rgb=NAVY_RGB)

    # ── 함소아 집행 금액 정산표 ──────────────────────────────
    s_bill = new_slide()
    aligo_mark(s_bill)
    section_title(s_bill, "6. 함소아 기사 집행 금액 정산표")

    if billing_data:
        bill_hdr = ['매체사', '분류', '개재 건수', '건당 견적', '매체별 합산']
        bill_table_data = [bill_hdr]
        total_amount = 0

        for rec in billing_data:
            vals = list(rec.values())
            row = (vals + [''] * 5)[:5]
            bill_table_data.append(row)
            try:
                amt_str = str(vals[4] if len(vals) > 4 else (vals[-1] if vals else '')).replace(',', '').replace('원', '').strip()
                if amt_str.isdigit():
                    total_amount += int(amt_str)
            except:
                pass

        bill_table_data.append(['TOTAL', f"기획기사 {len(billing_data) - 1}건 + 보고서 집계",
                                 f"{len(billing_data)}건", '', f"{total_amount:,}원" if total_amount else ''])

        add_table(s_bill, bill_table_data, [2.8, 2.5, 1.3, 1.3, 1.6], left=1.8, top=1.1, max_height=5.5)

        txt(s_bill, "※ 해당 금액은 VAT 별도 기준입니다",
            1.8, 1.1 + 0.55 * len(bill_table_data) + 0.15, 9.5, 0.35,
            size=9, italic=True, rgb=RGBColor(0x88, 0x88, 0x88))
    else:
        txt(s_bill, "정산 데이터가 없습니다.",
            0.5, 2.0, 12.0, 0.5, size=11, rgb=NAVY_RGB)

    # 저장
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ==================== 상세페이지 제작 관련 ====================

def analyze_product_for_detail_page(client, image_data_list, product_info, selling_points):
    """Claude Sonnet으로 제품 이미지+정보 분석, dict 반환"""
    content = []
    for img in image_data_list:
        content.append({
            "type": "image",
            "source": {"type": "base64", "media_type": img["media_type"], "data": img["data"]}
        })

    prompt = f"""제품 이미지와 정보를 분석하여 상세페이지 제작에 필요한 정보를 JSON 형식으로 반환해주세요.

[제품 정보]
{product_info}

[판매 포인트]
{selling_points if selling_points else "없음"}

다음 JSON 형식으로만 응답하세요 (설명 없이 JSON만):
{{
  "main_feature": "핵심 특징 한 줄 (20자 이내)",
  "features": ["특징1", "특징2", "특징3"],
  "mood": "모던/내추럴/프리미엄/귀여운 중 택1",
  "color_theme": "주요 색상 (예: 베이지, 딥블루)",
  "target": "타겟 고객 (예: 30대 주부, 직장인)",
  "tagline": "캐치프레이즈 한 줄 (25자 이내)"
}}"""

    content.append({"type": "text", "text": prompt})

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1000,
        messages=[{"role": "user", "content": content}]
    )

    raw = response.content[0].text.strip()
    try:
        if "```" in raw:
            raw = re.sub(r'```(?:json)?\n?', '', raw).strip().rstrip('`').strip()
        result = json.loads(raw)
    except Exception:
        result = {
            "main_feature": (product_info[:20] if product_info else "프리미엄 제품"),
            "features": ["고품질 소재", "편리한 사용", "합리적 가격"],
            "mood": "모던",
            "color_theme": "화이트",
            "target": "20-40대",
            "tagline": "당신의 일상을 특별하게"
        }
    return result


def generate_banner_copy(client, analysis, banner_type, product_name, selling_points, review_texts=None):
    """Claude Haiku로 배너별 카피라이팅 생성. banner_type: hero/features/usage/review/info"""

    type_context = {
        "hero": f"메인 히어로 배너. 핵심특징: {analysis.get('main_feature', '')} / 캐치프레이즈: {analysis.get('tagline', '')}",
        "features": f"제품 특징 배너. 특징: {', '.join(analysis.get('features', []))} / 판매포인트: {selling_points or '없음'}",
        "usage": f"사용법/활용 배너. 분위기: {analysis.get('mood', '')} / 타겟: {analysis.get('target', '')}",
        "review": f"구매 후기 배너. 참고 리뷰: {chr(10).join((review_texts or [])[:2]) or '없음'}",
        "info": f"브랜드/상품 정보 배너. 핵심특징: {analysis.get('main_feature', '')} / 분위기: {analysis.get('mood', '')}",
    }

    prompt = f"""상세페이지 배너 카피라이팅 작업입니다.

제품명: {product_name}
배너 유형: {banner_type}
배경 정보: {type_context.get(banner_type, '')}

다음 JSON 형식으로만 응답하세요 (설명 없이 JSON만):
{{
  "headline": "메인 헤드라인 (15자 이내)",
  "subtext": "서브 텍스트 (25자 이내)",
  "body": "본문 텍스트 (40자 이내)"
}}

[작성 조건]
- 한국어로 작성
- 구매욕구를 자극하는 카피
- 이모지 사용 금지"""

    response = client.messages.create(
        model="claude-haiku-4-5-20251001",
        max_tokens=300,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = response.content[0].text.strip()
    try:
        if "```" in raw:
            raw = re.sub(r'```(?:json)?\n?', '', raw).strip().rstrip('`').strip()
        result = json.loads(raw)
    except Exception:
        result = {
            "headline": product_name[:15],
            "subtext": analysis.get('tagline', '특별한 제품을 만나보세요')[:25],
            "body": (selling_points[:40] if selling_points else '지금 바로 확인해보세요')
        }
    return result


def create_banner_image(product_img_bytes, copy_data, banner_type, style, analysis,
                         font_regular_path, font_bold_path):
    """Pillow로 780x1200px 배너 이미지 생성. 반환: PIL Image"""
    try:
        from PIL import Image as _PILImg, ImageDraw as _Draw, ImageFont as _Font
    except ImportError:
        raise ImportError("Pillow 라이브러리가 필요합니다. pip install Pillow")

    W, H = 780, 1200
    IMG_H = int(H * 0.6)   # 상단 60% = 720px (제품 이미지 영역)
    TEXT_Y = IMG_H           # 텍스트 시작 Y = 720

    # ── 배경 생성 ──
    if style == "심플형":
        banner = _PILImg.new("RGB", (W, H), (255, 255, 255))
    else:
        # 감성형: 색상 테마 기반 그라디언트
        banner = _PILImg.new("RGB", (W, H), (255, 255, 255))
        draw_grad = _Draw(banner)
        color_theme = analysis.get('color_theme', '베이지').lower()
        if '블루' in color_theme or 'blue' in color_theme:
            top_c, bot_c = (220, 235, 255), (240, 248, 255)
        elif '핑크' in color_theme or 'pink' in color_theme:
            top_c, bot_c = (255, 230, 235), (255, 245, 248)
        elif '그린' in color_theme or 'green' in color_theme:
            top_c, bot_c = (220, 245, 230), (240, 255, 245)
        elif '퍼플' in color_theme or 'purple' in color_theme:
            top_c, bot_c = (240, 225, 255), (250, 240, 255)
        else:  # 기본 베이지/따뜻한 톤
            top_c, bot_c = (255, 248, 240), (250, 243, 233)
        for y in range(H):
            r = int(top_c[0] + (bot_c[0] - top_c[0]) * y / H)
            g = int(top_c[1] + (bot_c[1] - top_c[1]) * y / H)
            b = int(top_c[2] + (bot_c[2] - top_c[2]) * y / H)
            draw_grad.line([(0, y), (W, y)], fill=(r, g, b))

    draw = _Draw(banner)

    # ── 폰트 로드 (fallback 처리) ──
    def _load_font(path, size):
        try:
            if path and os.path.exists(path):
                return _Font.truetype(path, size)
        except Exception:
            pass
        try:
            return _Font.load_default()
        except Exception:
            return None

    font_title = _load_font(font_bold_path, 52)
    font_sub   = _load_font(font_regular_path, 32)
    font_body  = _load_font(font_regular_path, 26)
    font_small = _load_font(font_regular_path, 20)

    # ── 제품 이미지 배치 (상단 60%, 비율 유지, 중앙 정렬) ──
    if product_img_bytes:
        try:
            prod = _PILImg.open(io.BytesIO(product_img_bytes)).convert("RGBA")
            pad = 40
            prod.thumbnail((W - pad * 2, IMG_H - pad * 2), _PILImg.LANCZOS)
            pw, ph = prod.size
            px = (W - pw) // 2
            py = pad + ((IMG_H - pad * 2 - ph) // 2)
            if style == "심플형":
                bg_layer = _PILImg.new("RGBA", prod.size, (255, 255, 255, 255))
                composite = _PILImg.alpha_composite(bg_layer, prod).convert("RGB")
                banner.paste(composite, (px, py))
            else:
                banner.paste(prod, (px, py), prod)
        except Exception:
            pass

    # ── 구분선 (심플형) ──
    if style == "심플형":
        draw.line([(30, TEXT_Y), (W - 30, TEXT_Y)], fill=(220, 220, 220), width=1)

    # ── 텍스트 색상 ──
    if style == "심플형":
        text_color, sub_color, accent_color = (30, 30, 30), (80, 80, 80), (60, 60, 60)
    else:
        text_color, sub_color, accent_color = (40, 40, 40), (90, 80, 70), (120, 80, 50)

    headline = copy_data.get("headline", "")
    subtext  = copy_data.get("subtext", "")
    body     = copy_data.get("body", "")

    # ── 텍스트 그리기 헬퍼 ──
    def _draw_centered(text, y, font, color):
        if not text or not font:
            return y + 40
        try:
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            draw.text(((W - tw) // 2, y), text, font=font, fill=color)
            return y + th + 14
        except Exception:
            return y + 50

    def _draw_wrapped(text, y, font, color, max_w=700):
        if not text or not font:
            return y + 30
        try:
            line, result_y = "", y
            for ch in text:
                test = line + ch
                bbox = draw.textbbox((0, 0), test, font=font)
                if bbox[2] - bbox[0] > max_w and line:
                    result_y = _draw_centered(line, result_y, font, color)
                    line = ch
                else:
                    line = test
            if line:
                result_y = _draw_centered(line, result_y, font, color)
            return result_y
        except Exception:
            return y + 40

    # ── 배너 타입별 레이아웃 ──
    y = TEXT_Y + 35

    if banner_type == "hero":
        y = _draw_centered(headline, y, font_title, text_color)
        y += 12
        y = _draw_wrapped(subtext, y, font_sub, sub_color)
        y += 8
        y = _draw_wrapped(body, y, font_body, accent_color)

    elif banner_type == "features":
        y = _draw_centered(headline, y, font_sub, text_color)
        y += 22
        for feat in analysis.get('features', [])[:3]:
            y = _draw_wrapped(f"• {feat}", y, font_body, sub_color)
            y += 6

    elif banner_type == "usage":
        y = _draw_centered(headline, y, font_sub, text_color)
        y += 16
        y = _draw_wrapped(subtext, y, font_body, sub_color)
        y += 10
        y = _draw_wrapped(body, y, font_small, accent_color)

    elif banner_type == "review":
        _draw_centered('"', y, font_title, (180, 180, 180))
        y += 22
        y = _draw_wrapped(body, y, font_body, sub_color)
        y += 10
        y = _draw_centered(f"— {subtext}", y, font_small, accent_color)

    elif banner_type == "info":
        y = _draw_centered(headline, y, font_sub, text_color)
        y += 16
        y = _draw_wrapped(subtext, y, font_body, sub_color)
        y += 10
        y = _draw_wrapped(body, y, font_small, accent_color)

    return banner


def create_detail_zip(banners_a, banners_b, product_name):
    """A/B버전 배너 이미지 리스트를 ZIP으로 패키징. 반환: BytesIO"""
    zip_buf = io.BytesIO()
    labels = ["배너01", "배너02", "배너03", "배너04", "배너05"]

    with zipfile.ZipFile(zip_buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for i, img in enumerate(banners_a):
            if i >= len(labels):
                break
            img_buf = io.BytesIO()
            img.save(img_buf, format="JPEG", quality=92)
            zf.writestr(f"A버전_심플형/{product_name}_A_{labels[i]}.jpg", img_buf.getvalue())

        for i, img in enumerate(banners_b):
            if i >= len(labels):
                break
            img_buf = io.BytesIO()
            img.save(img_buf, format="JPEG", quality=92)
            zf.writestr(f"B버전_감성형/{product_name}_B_{labels[i]}.jpg", img_buf.getvalue())

    zip_buf.seek(0)
    return zip_buf


# ==================== Streamlit 앱 ====================

st.set_page_config(page_title="버즈필터 자동화", page_icon="🤖", layout="wide")

with st.sidebar:
    st.markdown("## 📋 메뉴")
    st.markdown("---")
    menu = st.radio("", options=[
        "🏭 버즈필터 발주",
        "🚚 위탁 발주",
        "✍️ 리뷰 생성",
        "📝 리뷰 입력",
        "📄 견적서 생성",
        "🖼️ 상세페이지 제작",
        "🌐 홈페이지 자동 개선",
        "📊 함소아 보고서",
        "🖼️ 배경 흰색 변환",
        "📖 바이럴 백과사전",
    ], label_visibility="collapsed")
    st.markdown("---")
    st.caption("버즈필터 업무 자동화 시스템")
    st.caption("© 2025 알리고미디어")

if menu == "🏭 버즈필터 발주":
    st.title("🤖 버즈필터 발주 자동 장부 입력")
    st.subheader("📊 발주서 엑셀을 업로드하면 장부에 자동으로 입력합니다.")
    uploaded_file = st.file_uploader("발주서 엑셀 파일 선택 (.xlsx)", type=['xlsx'])
    if uploaded_file:
        df = pd.read_excel(uploaded_file)
        df.columns = [col.split('(')[0].strip() for col in df.columns]
        st.write("📂 업로드 데이터 미리보기:", df.head())
        st.write(f"총 {len(df)}건 발주 데이터 확인")
        if st.button("🚀 장부 자동입력 시작"):
            client = get_anthropic_client()
            with st.spinner("마진계산기 불러오는 중..."):
                ms = get_sheet("2. 버즈필터 마진 계산기")
                if ms is None: st.stop()
                amd = ms.get_all_values()
                calc_df = pd.DataFrame(amd[2:], columns=amd[1])
                # 브랜드 컬럼 ffill: 첫 행에만 브랜드 있고 나머지 빈칸이므로 아래로 채움
                calc_df['브랜드'] = calc_df['브랜드'].replace('', pd.NA).ffill()
                calc_df = calc_df[calc_df['제품명'].str.strip() != '']
                st.success(f"✅ 마진계산기 로드 완료 ({len(calc_df)}개 상품)")
            with st.spinner("AI가 상품 매칭 중..."):
                today = datetime.now()
                rows_to_add, match_results = [], []
                product_col = '상품코드 표' if '상품코드 표' in calc_df.columns else '상품코드'

                # 전체 판매자코드 목록을 AI 프롬프트용으로 빌드 (루프 전 1회)
                catalog_lines = []
                for _, r in calc_df.iterrows():
                    brand = str(r.get('브랜드', '')).split('\n')[0].strip()
                    series = str(r.get('필터 시리즈', '')).strip()
                    product = str(r.get('제품명', '')).strip()
                    code = str(r.get(product_col, '')).strip()
                    if code and product:
                        catalog_lines.append(f"{code} | {brand} | {series} | {product}")
                catalog_str = '\n'.join(catalog_lines)

                qty_fail_items = []
                for idx, row in df.iterrows():
                    raw = str(row.get('상품명+옵션+개수', ''))
                    qty = extract_qty_from_text(raw)
                    if '/' not in raw:
                        qty_fail_items.append(raw)
                    ch = str(row.get('판매처', '쿠팡'))
                    price_raw = str(row.get('가격', '0')).replace(',', '').replace('원', '').strip()
                    try:
                        price = int(float(price_raw)) if price_raw else 0
                    except (ValueError, TypeError):
                        price = 0

                    prompt = f"""발주서 항목을 아래 판매자코드 전체 목록에서 찾아 매칭해줘.

[판매자코드 전체 목록]
코드 | 브랜드 | 필터시리즈 | 제품명
{catalog_str}

[발주서 항목]
{raw}

[매칭 우선순위]
1순위: 모델번호 숫자 정확 일치 (1245≠1645 → 숫자 다르면 무조건 다른 제품)
2순위: 필터시리즈 일치 (X툴≠Y툴, 3벌용≠5벌용, DH시리즈≠CDH시리즈 등)
2순위: 구성품 유사도 (헤파+탈취, 복합필터, 기능성 개수 등)
3순위: 가장 유사한 것

[규칙]
- 위 목록에서 반드시 하나 선택 (새 코드 생성 금지)
- 목록에 없을 때만: 상품코드: 미등록 / 브랜드: 미등록
- 마크다운/설명 출력 금지

[출력 형식 — 정확히 두 줄만]
상품코드: (목록의 코드)
브랜드: (목록의 브랜드)"""
                    try:
                        resp = client.messages.create(
                            model="claude-haiku-4-5-20251001", max_tokens=100,
                            messages=[{"role": "user", "content": prompt}]
                        )
                    except Exception as api_err:
                        st.error(f"❌ API 오류 ({type(api_err).__name__}): {api_err}")
                        st.stop()
                    rt = resp.content[0].text.strip()
                    mc, mb = parse_match_response(rt)
                    mb = str(mb).split('\n')[0].strip()
                    match_results.append({'상품명': raw, '매칭 브랜드': mb, '매칭 코드': mc, '판매처': ch, '가격': price, '수량(파싱)': qty})
                    rows_to_add.append([f"{today.year}년", f"{today.month}월", f"{today.day}일", mb, mc, ch, price, qty])
                st.session_state['rows_to_add'] = rows_to_add
                st.session_state['match_results'] = match_results
                st.session_state['qty_fail_items'] = qty_fail_items
                st.session_state['ready_to_insert'] = True
            st.success("✅ AI 매칭 완료!")
        if st.session_state.get('ready_to_insert'):
            rdf = pd.DataFrame(st.session_state['match_results'])
            st.write("🔍 AI 매칭 결과")
            st.dataframe(rdf)
            qty_fails = st.session_state.get('qty_fail_items', [])
            if qty_fails:
                with st.expander(f"⚠️ {len(qty_fails)}건 수량 파싱 불가 → 수량 1로 처리 (클릭해서 확인)"):
                    for item in qty_fails:
                        st.write(f"• {item}")
            unm = rdf[rdf['매칭 코드'] == '미등록']
            if len(unm) > 0:
                st.warning(f"⚠️ {len(unm)}건 상품 매칭 실패")
            if st.button("✅ 확인했습니다. 장부에 최종 입력합니다."):
                with st.spinner("장부 입력 중..."):
                    try:
                        ls = get_sheet("2. 버즈필터 장부")
                        if ls is None: st.stop()
                        sr = find_last_data_row(ls)
                        st.info(f"📍 {sr}행부터 입력 시작")
                        insert_row_safe(ls, sr, st.session_state['rows_to_add'])
                        st.success(f"🎉 총 {len(st.session_state['rows_to_add'])}건 입력 완료!")
                        st.balloons()
                        st.session_state['ready_to_insert'] = False
                        st.session_state['rows_to_add'] = []
                        st.session_state['match_results'] = []
                    except Exception as e:
                        st.error(f"❌ 입력 실패: {e}")

elif menu == "✍️ 리뷰 생성":
    st.title("✍️ AI 리뷰 자동 생성기")
    st.subheader("제품 정보를 입력하면 자연스럽고 다양한 리뷰를 생성해드립니다.")
    if 'generated_reviews' not in st.session_state:
        st.session_state.generated_reviews = []
    if 'review_edit_mode' not in st.session_state:
        st.session_state.review_edit_mode = False
    st.markdown("### 📦 STEP 1 — 제품 정보 입력")
    col_left, col_right = st.columns([2, 1])
    with col_left:
        product_info = st.text_area("제품 정보 (제품명, 카테고리, 특징 등)", height=150, placeholder="예)\n제품명: 콜라겐 마스크팩\n카테고리: 스킨케어\n특징: 저자극 성분, 수분 집중 케어")
        selling_points = st.text_area("소구점 / 강조할 내용 (선택)", height=100, placeholder="예) 피부 흡수력, 아침에 쓰기 좋음, 가성비")
    with col_right:
        product_images = st.file_uploader("제품 이미지 (선택, 여러 장 가능)", type=["jpg", "jpeg", "png", "webp"], accept_multiple_files=True)
        if product_images:
            for img in product_images:
                st.image(img, caption=img.name, use_container_width=True)
    st.markdown("### ⚙️ STEP 2 — 리뷰 설정")
    if 'review_batches' not in st.session_state:
        st.session_state.review_batches = [{'id': 0}]
        st.session_state._batch_next_id = 1
    for batch in list(st.session_state.review_batches):
        bid = batch['id']
        col1, col2, col3 = st.columns([4, 4, 1])
        with col1:
            st.number_input("리뷰 개수", min_value=1, max_value=200, value=st.session_state.get(f"rc_{bid}", 10), step=1, key=f"rc_{bid}")
        with col2:
            st.number_input("리뷰당 글자 수 (약)", min_value=50, max_value=500, value=st.session_state.get(f"cc_{bid}", 150), step=10, key=f"cc_{bid}")
        with col3:
            st.write(""); st.write("")
            if len(st.session_state.review_batches) > 1 and st.button("🗑️", key=f"del_{bid}", help="삭제"):
                st.session_state.review_batches = [b for b in st.session_state.review_batches if b['id'] != bid]
                st.rerun()
    if st.button("➕ 설정 추가"):
        new_id = st.session_state.get('_batch_next_id', len(st.session_state.review_batches))
        st.session_state._batch_next_id = new_id + 1
        st.session_state.review_batches.append({'id': new_id})
        st.rerun()
    st.markdown("---")
    total_reviews_sum = sum(st.session_state.get(f"rc_{b['id']}", 10) for b in st.session_state.review_batches)
    st.caption(f"💡 총 {total_reviews_sum}개 리뷰 생성 예정 ({len(st.session_state.review_batches)}개 설정)")
    if st.button("🚀 리뷰 생성 시작", type="primary", use_container_width=True):
        if not product_info.strip():
            st.error("❌ 제품 정보를 입력해주세요!")
        else:
            try:
                ai_client = get_anthropic_client()
                image_data_list = []
                media_type_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png", "webp": "image/webp"}
                if product_images:
                    for product_image in product_images:
                        product_image.seek(0)
                        img_bytes = product_image.read()
                        img_b64 = base64.b64encode(img_bytes).decode('utf-8')
                        ext = product_image.name.split('.')[-1].lower()
                        image_data_list.append({"media_type": media_type_map.get(ext, "image/jpeg"), "data": img_b64})
                batches_snapshot = list(st.session_state.review_batches)
                num_batches = len(batches_snapshot)
                progress_bar = st.progress(0)
                status_text = st.empty()
                all_reviews = []
                review_offset = 0
                for batch_idx, batch in enumerate(batches_snapshot):
                    bid = batch['id']
                    b_count = st.session_state.get(f"rc_{bid}", 10)
                    b_chars = st.session_state.get(f"cc_{bid}", 150)
                    status_text.text(f"🚀 설정 {batch_idx+1}/{num_batches} 생성 중 ({b_count}개, {b_chars}자)...")
                    def update_progress(current_batch, total_b, total_generated, _bidx=batch_idx, _nb=num_batches, _off=review_offset):
                        overall = (_bidx + current_batch / max(total_b, 1)) / _nb
                        progress_bar.progress(min(int(overall * 100), 99))
                        status_text.text(f"⏳ 설정 {_bidx+1}/{_nb} — 배치 {current_batch}/{total_b} — 현재까지 {_off + total_generated}개")
                    parsed = generate_reviews_with_claude(
                        client=ai_client, product_info=product_info, selling_points=selling_points,
                        review_count=b_count, char_count=b_chars,
                        image_data_list=image_data_list, progress_callback=update_progress)
                    for _, content in parsed:
                        review_offset += 1
                        all_reviews.append((review_offset, content))
                progress_bar.progress(100)
                status_text.empty()
                st.session_state.generated_reviews = all_reviews
                st.session_state.review_edit_mode = True
                st.success(f"✅ 총 {len(all_reviews)}개 리뷰 생성 완료!")
            except Exception as e:
                st.error(f"❌ 생성 실패: {e}")
    if st.session_state.review_edit_mode and st.session_state.generated_reviews:
        st.markdown("---")
        st.markdown(f"### 📋 STEP 3 — 결과 확인 및 수정 ({len(st.session_state.generated_reviews)}개)")
        updated_reviews = []
        for i, (num, content) in enumerate(st.session_state.generated_reviews):
            with st.expander(f"리뷰 {num}번", expanded=(i < 3)):
                edited = st.text_area(f"리뷰 {num} 내용", value=content, height=150, key=f"review_edit_{i}", label_visibility="collapsed")
                updated_reviews.append((num, edited))
        st.markdown("---")
        st.markdown("### 💾 STEP 4 — 저장 및 다운로드")
        col_save, col_reset = st.columns([3, 1])
        with col_save:
            if st.button("⬇️ 저장 및 엑셀 다운로드", type="primary", use_container_width=True):
                st.session_state.generated_reviews = updated_reviews
                excel_data = create_excel(updated_reviews)
                fname = f"리뷰_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
                st.download_button(label="📥 엑셀 파일 다운로드 클릭", data=excel_data, file_name=fname,
                                   mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                   use_container_width=True, type="primary")
                st.success("✅ 엑셀 파일이 준비되었습니다!")
        with col_reset:
            if st.button("🔄 초기화", use_container_width=True):
                st.session_state.generated_reviews = []
                st.session_state.review_edit_mode = False
                st.rerun()
        with st.expander("📄 텍스트 전체 보기 (복사용)"):
            full_text = ""
            for num, content in st.session_state.generated_reviews:
                full_text += f"{num}.\n{content}\n\n"
            st.text_area("전체 리뷰 텍스트", value=full_text.strip(), height=400, label_visibility="collapsed")

elif menu == "📝 리뷰 입력":
    st.title("📝 리뷰 엑셀 자동 변환기")
    st.subheader("리뷰 텍스트 파일을 업로드하면 엑셀 파일로 자동 변환합니다.")
    tab1, tab2 = st.tabs(["📁 파일 업로드", "✏️ 텍스트 직접 입력"])
    with tab1:
        utxt = st.file_uploader("리뷰 텍스트 파일 (.txt)", type=["txt"])
        if utxt:
            txt_content = utxt.read().decode("utf-8", errors="ignore")
            st.success(f"✅ {utxt.name} 업로드 완료")
            revs = parse_reviews(txt_content)
            if revs:
                st.markdown(f"### 📊 **{len(revs)}개** 리뷰 감지됨")
                with st.expander("👀 미리보기", expanded=True):
                    for num, content in revs[:5]:
                        st.markdown(f"**{num}번 리뷰**"); st.text(content[:200] + ("..." if len(content) > 200 else "")); st.divider()
                    if len(revs) > 5: st.info(f"... 외 {len(revs) - 5}개")
                st.download_button("⬇️ 엑셀 다운로드", create_excel(revs), "리뷰목록.xlsx", use_container_width=True, type="primary")
            else:
                st.error("❌ 리뷰를 파싱할 수 없습니다.")
    with tab2:
        mt = st.text_area("리뷰 내용 붙여넣기", height=300, placeholder='1\n리뷰 내용...\n\n2\n리뷰 내용...')
        if mt.strip():
            revs = parse_reviews(mt)
            if revs:
                st.markdown(f"### 📊 **{len(revs)}개** 리뷰 감지됨")
                with st.expander("👀 미리보기", expanded=True):
                    for num, content in revs[:5]:
                        st.markdown(f"**{num}번 리뷰**"); st.text(content[:200] + ("..." if len(content) > 200 else "")); st.divider()
                st.download_button("⬇️ 엑셀 다운로드", create_excel(revs), "리뷰목록.xlsx", use_container_width=True, type="primary")
            else:
                st.error("❌ 리뷰를 파싱할 수 없습니다.")

elif menu == "📄 견적서 생성":
    st.title("📄 견적서 자동 생성")
    st.subheader("정보를 입력하면 PDF 견적서를 자동으로 만들어드립니다.")
    col1, col2, col3 = st.columns(3)
    with col1: client_name = st.text_input("고객사명", placeholder="예) 지케이라이프")
    with col2: quote_date = st.text_input("견적일", value=datetime.now().strftime("%Y. %m. %d"))
    with col3: tax_type = st.radio("계산서 발행 여부", ["발행", "미발행"], horizontal=True)
    memo = st.text_input("비고 (선택)", placeholder="예) 패키지 할인 포함")
    st.markdown("---")
    st.markdown("### 📋 항목 입력")
    if 'quote_items' not in st.session_state:
        st.session_state.quote_items = [{"품목": "", "구성": "", "수량": 1, "단가": 0, "비고": ""}]
    hcols = st.columns([3, 2, 1, 2, 2, 1])
    for col, lbl in zip(hcols, ["**품목**", "**구성**", "**수량**", "**단가(원)**", "**공급가액**", "**삭제**"]):
        col.markdown(lbl)
    to_del = []
    for i, item in enumerate(st.session_state.quote_items):
        cols = st.columns([3, 2, 1, 2, 2, 1])
        item["품목"] = cols[0].text_input(f"p{i}", value=item["품목"], label_visibility="collapsed", placeholder="예) 쿠팡 리뷰")
        item["구성"] = cols[1].text_input(f"g{i}", value=item["구성"], label_visibility="collapsed", placeholder="예) 실행비")
        item["수량"] = cols[2].number_input(f"q{i}", value=item["수량"], min_value=0, label_visibility="collapsed")
        item["단가"] = cols[3].number_input(f"u{i}", value=item["단가"], min_value=0, step=100, label_visibility="collapsed")
        sp = item["수량"] * item["단가"]
        cols[4].markdown(f"<div style='padding:8px 0;font-weight:bold;'>₩{sp:,}</div>", unsafe_allow_html=True)
        if cols[5].button("🗑️", key=f"d{i}"): to_del.append(i)
    for i in sorted(to_del, reverse=True): st.session_state.quote_items.pop(i)
    if to_del: st.rerun()
    if st.button("➕ 항목 추가"):
        st.session_state.quote_items.append({"품목": "", "구성": "", "수량": 1, "단가": 0, "비고": ""})
        st.rerun()
    st.markdown("---")
    valid = [it for it in st.session_state.quote_items if it["품목"].strip()]
    sup = sum(it["수량"] * it["단가"] for it in valid)
    vat = int(sup * 0.1) if tax_type == "발행" else 0
    tot = sup + vat
    mc1, mc2, mc3 = st.columns(3)
    mc1.metric("공급가액 합계", f"₩{sup:,}")
    if tax_type == "발행":
        mc2.metric("세액 (VAT 10%)", f"₩{vat:,}")
        mc3.metric("최종 합계 (부가세 포함)", f"₩{tot:,}")
    else:
        mc2.metric("계산서", "미발행")
        mc3.metric("최종 합계 (VAT 없음)", f"₩{tot:,}")
    st.markdown("---")
    if st.button("📄 견적서 PDF 생성", type="primary", use_container_width=True):
        if not client_name.strip():
            st.error("❌ 고객사명을 입력해주세요!")
        elif not valid:
            st.error("❌ 항목을 최소 1개 이상 입력해주세요!")
        else:
            with st.spinner("PDF 생성 중..."):
                BASE_DIR = os.path.dirname(os.path.abspath(__file__))
                stamp_path = os.path.join(BASE_DIR, '직인_투명.png')
                qd = {"date": quote_date, "client": client_name, "tax_type": tax_type, "memo": memo, "items": valid}
                try:
                    pdf_buf = generate_quote_pdf(qd, stamp_path)
                    fname = f"견적서_{client_name}_{quote_date.replace('. ', '').replace('.', '')}.pdf"
                    st.success("✅ 견적서 PDF 생성 완료!")
                    st.download_button("⬇️ PDF 다운로드", pdf_buf, fname, mime="application/pdf", use_container_width=True, type="primary")
                except Exception as e:
                    st.error(f"❌ PDF 생성 실패: {e}")

elif menu == "🌐 홈페이지 자동 개선":
    st.title("🌐 홈페이지 자동 개선 + 자동 배포")
    st.subheader("HTML과 이미지를 업로드하면 Claude가 수정하고 Netlify에 자동 배포합니다.")
    try:
        NETLIFY_TOKEN = st.secrets["NETLIFY_TOKEN"]
        NETLIFY_SITE_ID = st.secrets["NETLIFY_SITE_ID"]
        netlify_ready = True
        st.success("✅ Netlify 연결 준비 완료 — 버튼 한 번으로 자동 배포됩니다.")
    except:
        netlify_ready = False
        st.error("❌ Netlify 미연결 — secrets.toml에 NETLIFY_TOKEN과 NETLIFY_SITE_ID를 추가해주세요.")
    st.markdown("---")
    st.markdown("### 📂 STEP 1 — index.html 업로드")
    uploaded_html = st.file_uploader("index.html 업로드", type=["html", "htm"], key="html_upload")
    st.markdown("### 🖼️ STEP 2 — 이미지 파일 업로드")
    uploaded_images = st.file_uploader("이미지 파일 (여러 개 동시 선택 가능)", type=["png", "jpg", "jpeg", "gif", "webp", "ico"], accept_multiple_files=True, key="image_upload")
    if uploaded_images:
        st.success(f"✅ 이미지 {len(uploaded_images)}개: {', '.join([f.name for f in uploaded_images])}")
    st.markdown("---")
    st.markdown("### ✅ STEP 3 — 개선 항목 선택")
    col1, col2, col3 = st.columns(3)
    with col1: check_mobile = st.checkbox("📱 모바일 최적화", value=True)
    with col2: check_responsive = st.checkbox("📐 반응형 디자인", value=True)
    with col3: check_seo = st.checkbox("🔍 구글 SEO", value=True)
    check_extra = st.text_area("📝 추가 요청사항 (선택)", placeholder="예) 버튼 색상을 더 눈에 띄게", height=80)
    st.markdown("---")
    if uploaded_html:
        html_content = uploaded_html.read().decode("utf-8", errors="ignore")
        if not any([check_mobile, check_responsive, check_seo, check_extra.strip()]):
            st.warning("⚠️ 개선 항목을 최소 1개 이상 선택해주세요.")
        else:
            if st.button("🚀 Claude 수정 + Netlify 자동 배포", type="primary", use_container_width=True):
                check_list = []
                if check_mobile: check_list.append("1. 모바일 최적화")
                if check_responsive: check_list.append("2. 반응형 디자인")
                if check_seo: check_list.append("3. 구글 SEO")
                if check_extra.strip(): check_list.append(f"4. 추가 요청: {check_extra.strip()}")
                prompt = f"""너는 웹 개발 전문가야. 아래 HTML을 분석하고 수정해서 완성된 HTML 코드만 반환해줘.
[개선 항목]
{chr(10).join(check_list)}
[주의사항]
- 수정된 HTML 전체 코드만 반환 (설명 없이 <!DOCTYPE html>부터 시작)
- 기존 디자인·색상·브랜드 정체성 유지
- 한국어 텍스트 수정 금지
[원본 HTML]
{html_content}"""
                progress_bar = st.progress(0)
                status_text = st.empty()
                try:
                    status_text.text("🤖 Claude가 분석 및 수정 중...")
                    progress_bar.progress(20)
                    ai_client = get_anthropic_client()
                    response = ai_client.messages.create(model="claude-opus-4-5", max_tokens=16000, messages=[{"role": "user", "content": prompt}])
                    improved_html = response.content[0].text.strip()
                    if improved_html.startswith("```"):
                        lines = improved_html.split("\n")
                        improved_html = "\n".join(lines[1:-1] if lines[-1].strip() == "```" else lines[1:])
                    progress_bar.progress(60)
                    status_text.text("✅ Claude 수정 완료! Netlify 배포 중...")
                    extra_files = {}
                    if uploaded_images:
                        for img_file in uploaded_images:
                            img_file.seek(0)
                            extra_files[img_file.name] = img_file.read()
                    if netlify_ready:
                        success, result = deploy_to_netlify(improved_html, NETLIFY_SITE_ID, NETLIFY_TOKEN, extra_files)
                        progress_bar.progress(100)
                        if success:
                            status_text.text("🎉 완료!")
                            st.success("🎉 수정 완료 + Netlify 자동 배포 성공!")
                            st.balloons()
                            st.link_button("🔗 aligomedia.co.kr 확인하기", "https://aligomedia.co.kr")
                        else:
                            st.error(f"❌ Netlify 배포 실패\n\n{result}")
                    else:
                        progress_bar.progress(100)
                        status_text.text("✅ 수정 완료")
                    st.session_state["improved_html"] = improved_html
                    st.session_state["original_html"] = html_content
                    st.session_state["improvement_done"] = True
                except Exception as e:
                    st.error(f"❌ 오류 발생: {e}")
    else:
        st.info("👆 STEP 1에서 index.html을 업로드하면 시작할 수 있어요!")
    if st.session_state.get("improvement_done"):
        improved_html = st.session_state["improved_html"]
        original_html = st.session_state["original_html"]
        st.markdown("---")
        col_before, col_after = st.columns(2)
        with col_before:
            st.markdown("#### 📄 수정 전")
            st.metric("파일 크기", f"{len(original_html):,}자")
            with st.expander("원본 코드 보기"):
                st.code(original_html[:1500] + "...", language="html")
        with col_after:
            st.markdown("#### ✅ 수정 후")
            st.metric("파일 크기", f"{len(improved_html):,}자", delta=f"{len(improved_html) - len(original_html):+,}자")
            with st.expander("수정된 코드 보기"):
                st.code(improved_html[:1500] + "...", language="html")
        st.markdown("---")
        st.download_button(label="⬇️ 수정된 index.html 다운로드", data=improved_html.encode("utf-8"), file_name="index.html", mime="text/html", use_container_width=True)
        if st.button("🔄 처음부터 다시", use_container_width=True):
            st.session_state["improvement_done"] = False
            st.session_state["improved_html"] = ""
            st.rerun()

elif menu == "📊 함소아 보고서":
    st.title("📊 함소아한의원 보고서 분석 텍스트 생성")

    # ── STEP 1: 기본 설정 ──
    st.markdown("### STEP 1. 기본 설정")
    col1, col2 = st.columns(2)
    with col1:
        report_date = st.text_input("보고서 기준일",
                                    value=datetime.now().strftime("%Y년 %m월 %d일"),
                                    placeholder="예) 2026년 3월 20일")
        comp_sheet_name = st.text_input("경쟁사 시트 탭 이름", value="경쟁사 동향 분석 시트")
    with col2:
        report_month = st.text_input("보고 월",
                                     value=datetime.now().strftime("%Y년 %m월"),
                                     placeholder="예) 2026년 3월")
        mgmt_sheet_name = st.text_input("함소아 관리 시트 탭 이름", value="함소아 한의원 관리 시트")

    hamsoa_manual = st.number_input(
        "함소아 기사 발행 건수 (0이면 시트에서 자동 계산)",
        min_value=0, value=0)

    st.markdown("---")

    # ── STEP 2: 자생한방병원 전략 데이터 입력 ──
    st.markdown("### STEP 2. 자생한방병원 전략 데이터 입력")
    st.caption("전략유형 이름을 먼저 쓰고, 그 아래에 '키워드 건수' 형식으로 입력. 건수 없는 줄은 유형 헤더로 처리됩니다.")

    jasaeng_raw = st.text_area(
        "자생한방병원 전략 데이터",
        height=300,
        placeholder="""질환 타겟형
비염 1
일자목 1
혈당스파이크 24

시술 중심형

마케팅형
척추건강 증진사업 21
공동모금회 백미 29

브랜드 강화형
동작침법 52
발목 염좌 18""",
        help="숫자가 있는 줄 = 항목, 숫자 없는 줄 = 전략유형 헤더"
    )

    st.markdown("---")

    # ── STEP 3: 생성 버튼 ──
    if st.button("✨ 분석 텍스트 생성", type="primary", use_container_width=True):

        if not jasaeng_raw.strip():
            st.error("자생한방병원 전략 데이터를 입력해주세요.")
            st.stop()

        jasaeng_strategies, jasaeng_order, jasaeng_totals, jasaeng_grand_total = \
            parse_jasaeng_strategy(jasaeng_raw)

        with st.spinner("경쟁사 시트 불러오는 중..."):
            comp_sheet = get_hamsoa_sheet(comp_sheet_name)
            if comp_sheet is None:
                st.error(f"'{comp_sheet_name}' 시트를 찾을 수 없습니다.")
                st.stop()
            competitor_records, meta = parse_competitor_sheet(comp_sheet)

        with st.spinner("함소아 관리 시트 불러오는 중..."):
            mgmt_sheet = get_hamsoa_sheet(mgmt_sheet_name)
            if mgmt_sheet is None:
                st.error(f"'{mgmt_sheet_name}' 시트를 찾을 수 없습니다.")
                st.stop()
            hamsoa_articles, billing_data = parse_hamsoa_sheet(mgmt_sheet)

        # 함소아 기사 수: 수동 입력 우선, 없으면 시트 전체 행 수
        if hamsoa_manual > 0:
            hamsoa_count = hamsoa_manual
        else:
            hamsoa_count = len(hamsoa_articles)
        hospital_counts = meta.get('hospital_counts', {})

        st.success(f"✅ 자생 {jasaeng_grand_total}건 | 경쟁사 기사 {len(competitor_records)}건 | 함소아 {hamsoa_count}건 로드 완료")

        # 함소아 시트 파싱 결과 디버그
        if hamsoa_articles:
            with st.expander(f"📋 함소아 기사 현황 ({hamsoa_count}건) — 컬럼 확인", expanded=False):
                df_art = pd.DataFrame(hamsoa_articles)
                st.write("감지된 컬럼:", list(df_art.columns))
                # 진행현황 컬럼 찾아서 분포 보여주기
                status_col = next((c for c in df_art.columns if '진행' in c or '현황' in c), None)
                if status_col:
                    st.write(f"'{status_col}' 값 분포:", df_art[status_col].value_counts().to_dict())
                st.dataframe(df_art.head(10))
        else:
            st.warning("⚠️ 함소아 기사를 시트에서 읽지 못했습니다. '발행일', '구분', '제목' 컬럼이 있는 헤더 행이 있는지 확인해주세요.")

        # 자생 전략 집계 확인
        with st.expander("📊 자생한방병원 전략 집계 확인", expanded=True):
            cols = st.columns(len(jasaeng_order) if jasaeng_order else 1)
            for i, s in enumerate(jasaeng_order):
                t = jasaeng_totals.get(s, 0)
                items = jasaeng_strategies.get(s, [])
                with cols[i % len(cols)]:
                    st.metric(s, f"{t}건")
                    if items:
                        st.caption(" / ".join(f"{kw}({cnt})" for kw, cnt in items))
            st.markdown(f"**총 합계: {jasaeng_grand_total}건**")

        # ── AI 분석 텍스트 생성 ──
        st.markdown("---")
        st.markdown("### 🤖 AI 전문가 분석 텍스트")
        st.caption("데이터를 바탕으로 AI가 전문 분석가 어조로 작성합니다.")

        if st.button("✨ AI로 분석 텍스트 생성", type="primary", use_container_width=True):
            with st.spinner("AI가 분석 중... (5~10초)"):
                try:
                    ai_client = get_anthropic_client()
                    ai_text = generate_ai_report_text(
                        ai_client,
                        jasaeng_totals, jasaeng_order, jasaeng_grand_total,
                        jasaeng_strategies, hospital_counts, hamsoa_count,
                        competitor_records, report_date
                    )
                    st.session_state["ai_report_text"] = ai_text
                except Exception as e:
                    st.error(f"AI 생성 실패: {e}")

        if st.session_state.get("ai_report_text"):
            raw = st.session_state["ai_report_text"]
            # 섹션1, 섹션2 분리 표시
            if "섹션1" in raw or "섹션2" in raw:
                parts = re.split(r'={2,}\s*섹션[12]\s*={2,}', raw)
                labels = ["섹션1: 발행 수량 비교", "섹션2: 전략유형별 세부 분석"]
                sections = [p.strip() for p in parts if p.strip()]
                for i, sec in enumerate(sections):
                    label = labels[i] if i < len(labels) else f"섹션{i+1}"
                    st.markdown(f"**📝 {label}**")
                    st.text_area("복사해서 PPT에 붙여넣으세요", value=sec, height=220, key=f"ai_sec_{i}")
            else:
                st.text_area("복사해서 PPT에 붙여넣으세요", value=raw, height=420, key="ai_full")

        st.markdown("---")

        # ── 기본 템플릿 텍스트 (참고용) ──
        with st.expander("📋 기본 템플릿 텍스트 (참고용)", expanded=False):
            qty_texts = generate_quantity_analysis(
                hospital_counts, hamsoa_count, jasaeng_grand_total, report_date)
            qty_output = "\n\n".join(f"➡  {t}" for t in qty_texts)
            st.text_area("수량 비교", value=qty_output, height=180, key="qty_out")

            strat_texts = generate_strategy_analysis(
                jasaeng_strategies, jasaeng_order, jasaeng_totals, jasaeng_grand_total,
                competitor_records, hospital_counts)
            strat_output = "\n\n".join(f"➡  {t}" for t in strat_texts)
            st.text_area("전략유형 분석", value=strat_output, height=200, key="strat_out")

        # ── 데이터 미리보기 ──
        with st.expander("📋 경쟁사 기사 미리보기", expanded=False):
            if competitor_records:
                st.dataframe(pd.DataFrame(competitor_records).head(20))
                if hospital_counts:
                    st.write("병원별 건수:", hospital_counts)
            else:
                st.warning("⚠️ 경쟁사 데이터 없음")

        with st.expander("📋 함소아 기사 미리보기", expanded=False):
            if hamsoa_articles:
                st.dataframe(pd.DataFrame(hamsoa_articles))
            else:
                st.warning("⚠️ 함소아 기사 데이터 없음")

elif menu == "🚚 위탁 발주":
    import io as _io, re as _re
    st.title("\U0001f69a 위탁 발주서 자동 생성")
    st.subheader("쿠팡 / 스마트스토어 주문 엑셀을 업로드하면 발주서를 자동으로 만들어 드립니다.")

    COLS_OUT = [
        "수령자명", "전화번호1", "전화번호2", "수령주소",
        "상품명+옵션+개수 (최대한 요약해서 작성 부탁드립니다)",
        "배송요청사항", "가격", "판매처"
    ]

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("#### \U0001f6d2 쿠팡 주문 엑셀")
        cp_file = st.file_uploader("쿠팡 발주 엑셀 업로드", type=["xlsx", "xls"], key="cp_upload")
    with col2:
        st.markdown("#### \U0001f3ea 스마트스토어 주문 엑셀")
        ss_file = st.file_uploader("스마트스토어 발주 엑셀 업로드", type=["xlsx", "xls"], key="ss_upload")
        ss_pw = st.text_input("암호 (없으면 비워두세요)", type="password", key="ss_pw")

    if st.button("\U0001f4cb 발주서 생성", type="primary"):
        if not cp_file and not ss_file:
            st.warning("\u26a0\ufe0f 쿠팡 또는 스마트스토어 파일을 하나 이상 업로드해주세요.")
        else:
            all_rows = []

            # 쿠팡 파싱
            if cp_file:
                try:
                    df_cp = pd.read_excel(cp_file, header=0, dtype=str)
                    cnt = 0
                    for _, row in df_cp.iterrows():
                        product = str(row.get("등록상품명", "") or "").strip()
                        option  = str(row.get("등록옵션명", "") or "").strip()
                        qty_s   = str(row.get("구매수(수량)", "1") or "1").strip()
                        qty = int(float(qty_s)) if qty_s.replace(".", "").isdigit() else 1
                        option_clean = _re.sub(r"^\d+개\s*", "", option).strip()
                        pstr = f"{product}, {option_clean} / {qty}개" if option_clean else f"{product} / {qty}개"
                        tel1 = str(row.get("수취인전화번호", "") or "").strip()
                        tel2 = str(row.get("구매자전화번호", "") or "").strip()
                        price_s = str(row.get("결제액", "0") or "0").replace(",", "").strip()
                        price = int(float(price_s)) if price_s.replace(".", "").isdigit() else 0
                        all_rows.append({
                            "수령자명": str(row.get("수취인이름", "") or "").strip(),
                            "전화번호1": tel1,
                            "전화번호2": tel2 if tel2 != tel1 else "",
                            "수령주소": str(row.get("수취인 주소", "") or "").strip(),
                            "상품명+옵션+개수 (최대한 요약해서 작성 부탁드립니다)": pstr,
                            "배송요청사항": str(row.get("배송메세지", "") or "").strip(),
                            "가격": price,
                            "판매처": "쿠팡"
                        })
                        cnt += 1
                    st.success(f"\u2705 쿠팡 {cnt}건 파싱 완료")
                except Exception as e:
                    st.error(f"\u274c 쿠팡 파싱 오류: {e}")

            # 스마트스토어 파싱
            if ss_file:
                try:
                    file_bytes = ss_file.read()
                    if ss_pw:
                        import msoffcrypto as _msoff
                        of = _msoff.OfficeFile(_io.BytesIO(file_bytes))
                        of.load_key(password=ss_pw)
                        target = _io.BytesIO()
                        of.decrypt(target)
                        target.seek(0)
                    else:
                        target = _io.BytesIO(file_bytes)
                    df_ss = pd.read_excel(target, header=1, dtype=str)
                    if "판매자 상품코드" not in df_ss.columns:
                        target.seek(0)
                        df_ss = pd.read_excel(target, header=0, dtype=str)
                    cnt = 0
                    for _, row in df_ss.iterrows():
                        code  = str(row.get("판매자 상품코드", "") or "").strip()
                        opt   = str(row.get("옵션정보", "") or "").strip()
                        qty_s = str(row.get("수량", "1") or "1").strip()
                        qty = int(float(qty_s)) if qty_s.replace(".", "").isdigit() else 1
                        if ": " in opt:
                            opt = opt.split(": ", 1)[1].strip()
                        if opt in ("nan", "None", ""):
                            opt = ""
                        pstr = f"{code}, {opt} / {qty}개" if opt else f"{code} / {qty}개"
                        addr1 = str(row.get("기본배송지", "") or "").strip()
                        addr2 = str(row.get("상세배송지", "") or "").strip()
                        address = f"{addr1} {addr2}".strip()
                        price_s = str(row.get("최종 상품별 주문금액", "0") or "0").replace(",", "").strip()
                        price = int(float(price_s)) if price_s.replace(".", "").isdigit() else 0
                        tel2 = str(row.get("수취인연락처2", "") or "").strip()
                        if tel2 in ("nan", "None"): tel2 = ""
                        all_rows.append({
                            "수령자명": str(row.get("수취인명", "") or "").strip(),
                            "전화번호1": str(row.get("수취인연락처1", "") or "").strip(),
                            "전화번호2": tel2,
                            "수령주소": address,
                            "상품명+옵션+개수 (최대한 요약해서 작성 부탁드립니다)": pstr,
                            "배송요청사항": str(row.get("배송메세지", "") or "").strip(),
                            "가격": price,
                            "판매처": "스마트스토어"
                        })
                        cnt += 1
                    st.success(f"\u2705 스마트스토어 {cnt}건 파싱 완료")
                except Exception as e:
                    st.error(f"\u274c 스마트스토어 파싱 오류: {e}")

            if all_rows:
                st.session_state["위탁발주_df"] = pd.DataFrame(all_rows, columns=COLS_OUT)
            else:
                st.warning("\u26a0\ufe0f 파싱된 데이터가 없습니다. 파일 형식을 확인해주세요.")

    # 결과 미리보기 + 다운로드
    if st.session_state.get("위탁발주_df") is not None:
        result_df = st.session_state["위탁발주_df"]
        st.markdown(f"### \U0001f4cb 발주서 미리보기 — 총 {len(result_df)}건")
        st.dataframe(result_df, use_container_width=True)
        today = datetime.now()
        fname = f"발주서_버즈필터 {today.month}월 {today.day}일.xlsx"
        buf = _io.BytesIO()
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            result_df.to_excel(writer, index=False, sheet_name="발주발송관리")
        st.download_button(
            label=f"\u2b07\ufe0f {fname} 다운로드",
            data=buf.getvalue(),
            file_name=fname,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            type="primary"
        )

elif menu == "🖼️ 배경 흰색 변환":
    st.title("🖼️ 배경 흰색 변환")
    st.subheader("투명 배경 이미지를 흰색 배경으로 자동 변환합니다.")

    from PIL import Image as PILImage

    uploaded = st.file_uploader("이미지 업로드 (PNG, JPG, WEBP)", type=["png", "jpg", "jpeg", "webp"])

    if uploaded:
        img = PILImage.open(uploaded)

        # 흰 배경 합성
        white_bg = PILImage.new("RGB", img.size, (255, 255, 255))
        if img.mode in ("RGBA", "LA"):
            white_bg.paste(img, mask=img.split()[-1])
        elif img.mode == "P" and "transparency" in img.info:
            white_bg.paste(img.convert("RGBA"), mask=img.convert("RGBA").split()[-1])
        else:
            white_bg.paste(img)

        col1, col2 = st.columns(2)
        with col1:
            st.markdown("**원본**")
            st.image(img, use_container_width=True)
        with col2:
            st.markdown("**변환 결과 (흰 배경)**")
            st.image(white_bg, use_container_width=True)

        # 다운로드 (JPG)
        buf = io.BytesIO()
        white_bg.save(buf, format="JPEG", quality=95)
        buf.seek(0)
        original_name = uploaded.name.rsplit(".", 1)[0]
        st.download_button(
            label="⬇️ 흰 배경 이미지 다운로드 (JPG)",
            data=buf,
            file_name=f"{original_name}_white.jpg",
            mime="image/jpeg",
            type="primary",
            use_container_width=True,
        )

elif menu == "🖼️ 상세페이지 제작":
    st.title("🖼️ 상세페이지 자동 제작")
    st.subheader("제품 이미지와 정보를 입력하면 A/B 버전 상세페이지 배너를 자동 생성합니다.")

    # session_state 초기화 (dp_ 프리픽스)
    for _dp_key, _dp_default in [
        ("dp_analysis", None), ("dp_copy_list", []),
        ("dp_banners_a", []), ("dp_banners_b", []),
        ("dp_reviews_used", [])
    ]:
        if _dp_key not in st.session_state:
            st.session_state[_dp_key] = _dp_default

    # ── STEP 1: 상품 정보 + 이미지 ──
    st.markdown("### STEP 1 — 상품 정보 + 이미지")

    dp_product_name = st.text_input(
        "상품명", placeholder="예) 유기농 콜라겐 마스크팩", key="dp_product_name_input")
    dp_selling_points = st.text_area(
        "판매 포인트 (선택)", height=80,
        placeholder="예) 저자극 성분, 수분 집중 케어, 가성비",
        key="dp_selling_points_input")

    dp_images = st.file_uploader(
        "제품 이미지 (최대 5장, jpg/png/webp)",
        type=["jpg", "jpeg", "png", "webp"],
        accept_multiple_files=True,
        key="dp_images_upload")

    if dp_images and len(dp_images) > 5:
        st.warning("최대 5장까지만 사용됩니다. 앞의 5장이 적용됩니다.")
        dp_images = list(dp_images)[:5]

    # 리뷰 불러오기
    dp_use_reviews = st.checkbox(
        "리뷰 생성 메뉴에서 만든 리뷰 불러오기", key="dp_use_reviews_chk")
    if dp_use_reviews:
        existing_reviews = st.session_state.get("generated_reviews", [])
        if existing_reviews:
            top_reviews = existing_reviews[:3]
            st.info(f"리뷰 {len(existing_reviews)}개 중 상위 3개를 배너에 활용합니다.")
            for _num, _content in top_reviews:
                st.caption(f"• {_content[:60]}...")
            st.session_state["dp_reviews_used"] = [c for _, c in top_reviews]
        else:
            st.warning("생성된 리뷰가 없습니다. 리뷰 생성 메뉴에서 먼저 리뷰를 만들어주세요.")
            st.session_state["dp_reviews_used"] = []
    else:
        st.session_state["dp_reviews_used"] = []

    if st.button("🔍 AI 상품 분석 시작", key="dp_analyze_btn"):
        if not dp_product_name.strip():
            st.error("상품명을 입력해주세요!")
        else:
            with st.spinner("Claude Sonnet이 상품을 분석 중..."):
                try:
                    _dp_client = get_anthropic_client()
                    _dp_img_list = []
                    _mt_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
                               "png": "image/png", "webp": "image/webp"}
                    if dp_images:
                        for _f in list(dp_images)[:5]:
                            _f.seek(0)
                            _b64 = base64.b64encode(_f.read()).decode('utf-8')
                            _ext = _f.name.split('.')[-1].lower()
                            _dp_img_list.append({
                                "media_type": _mt_map.get(_ext, "image/jpeg"),
                                "data": _b64
                            })
                    _analysis = analyze_product_for_detail_page(
                        _dp_client, _dp_img_list, dp_product_name, dp_selling_points)
                    st.session_state["dp_analysis"] = _analysis
                    st.success("분석 완료! 아래에서 결과를 확인하세요.")
                except Exception as _e:
                    st.error(f"분석 실패: {_e}")

    # 분석 결과 표시
    if st.session_state.get("dp_analysis"):
        _a = st.session_state["dp_analysis"]
        with st.expander("분석 결과 확인", expanded=True):
            _c1, _c2, _c3 = st.columns(3)
            with _c1:
                st.metric("분위기", _a.get("mood", ""))
                st.metric("색상 테마", _a.get("color_theme", ""))
            with _c2:
                st.metric("타겟", _a.get("target", ""))
                st.metric("핵심 특징", _a.get("main_feature", ""))
            with _c3:
                st.metric("캐치프레이즈", _a.get("tagline", ""))
            if _a.get("features"):
                st.write("감지된 특징:", " / ".join(_a.get("features", [])))

    # ── STEP 2: 배너 설정 (분석 완료 후만 표시) ──
    if st.session_state.get("dp_analysis"):
        st.markdown("---")
        st.markdown("### STEP 2 — 배너 설정")

        _scol1, _scol2 = st.columns(2)
        with _scol1:
            dp_banner_count = st.selectbox(
                "배너 장수", ["5장", "3장"], key="dp_banner_count_sel")
        with _scol2:
            dp_competitor_url = st.text_input(
                "경쟁사 URL (선택 — 차별화 포인트 분석용)",
                placeholder="https://example.com/product",
                key="dp_comp_url_input")

        if dp_competitor_url.strip():
            if st.button("경쟁사 분석", key="dp_comp_analyze_btn"):
                with st.spinner("경쟁사 페이지 분석 중..."):
                    try:
                        _headers = {
                            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
                                          'AppleWebKit/537.36 (KHTML, like Gecko) '
                                          'Chrome/120.0.0.0 Safari/537.36'
                        }
                        _resp = requests.get(
                            dp_competitor_url.strip(), headers=_headers, timeout=10)
                        _resp.encoding = 'utf-8'
                        _html_text = _resp.text[:5000]

                        _dp_client2 = get_anthropic_client()
                        _comp_prompt = f"""다음은 경쟁사 페이지 HTML입니다.
이 상세페이지와 차별화할 수 있는 포인트 3가지를 간결하게 분석해주세요.

HTML:
{_html_text}

번호 목록 형식으로 3가지만 작성해주세요."""
                        _cr = _dp_client2.messages.create(
                            model="claude-haiku-4-5-20251001",
                            max_tokens=500,
                            messages=[{"role": "user", "content": _comp_prompt}]
                        )
                        st.info(_cr.content[0].text.strip())
                    except requests.exceptions.Timeout:
                        st.warning("URL 연결 시간이 초과되었습니다. (10초)")
                    except requests.exceptions.RequestException as _re_err:
                        st.warning(f"URL 접근 실패: {_re_err}")
                    except Exception as _ce:
                        st.warning(f"분석 중 오류가 발생했습니다: {_ce}")

        # ── STEP 3: 생성 ──
        st.markdown("---")
        st.markdown("### STEP 3 — 배너 생성")

        if st.button("✨ 상세페이지 A/B 버전 생성", type="primary",
                     use_container_width=True, key="dp_generate_btn"):
            if not dp_images:
                st.info("제품 이미지가 없어도 생성됩니다. 이미지를 업로드하면 더 좋은 결과를 얻을 수 있습니다.")

            _dp_analysis = st.session_state["dp_analysis"]
            _dp_reviews  = st.session_state.get("dp_reviews_used", [])
            _banner_n    = 5 if dp_banner_count == "5장" else 3

            # 배너 타입 순서 결정
            if _dp_reviews:
                _btypes_all = ["hero", "features", "usage", "review", "info"]
            else:
                _btypes_all = ["hero", "features", "usage", "info", "info"]
            _btypes = _btypes_all[:_banner_n]

            # 폰트 경로 (generate_quote_pdf와 동일 방식)
            _dp_base = os.path.dirname(os.path.abspath(__file__))
            _font_r = os.path.join(_dp_base, 'NotoSansKR-Regular.ttf')
            _font_b = os.path.join(_dp_base, 'NotoSansKR-Bold.ttf')

            # 제품 이미지 바이트 (첫 번째 이미지)
            _prod_bytes = None
            if dp_images:
                try:
                    list(dp_images)[0].seek(0)
                    _prod_bytes = list(dp_images)[0].read()
                except Exception:
                    pass

            _dp_client3 = get_anthropic_client()
            _banners_a, _banners_b, _copy_list = [], [], []

            _progress = st.progress(0)
            _status   = st.empty()
            _total_steps = _banner_n * 2  # 카피 + 이미지 생성

            try:
                _step = 0
                for _i, _btype in enumerate(_btypes):
                    _status.text(f"배너 {_i+1}/{_banner_n} 카피 생성 중...")
                    _copy = generate_banner_copy(
                        _dp_client3, _dp_analysis, _btype,
                        dp_product_name, dp_selling_points, _dp_reviews)
                    _copy_list.append(_copy)
                    _step += 1
                    _progress.progress(int(_step / _total_steps * 100))

                    _status.text(f"배너 {_i+1}/{_banner_n} 이미지 생성 중 (A/B)...")
                    _ban_a = create_banner_image(
                        _prod_bytes, _copy, _btype, "심플형",
                        _dp_analysis, _font_r, _font_b)
                    _ban_b = create_banner_image(
                        _prod_bytes, _copy, _btype, "감성형",
                        _dp_analysis, _font_r, _font_b)
                    _banners_a.append(_ban_a)
                    _banners_b.append(_ban_b)
                    _step += 1
                    _progress.progress(int(_step / _total_steps * 100))

                st.session_state["dp_banners_a"] = _banners_a
                st.session_state["dp_banners_b"] = _banners_b
                st.session_state["dp_copy_list"] = _copy_list
                _progress.progress(100)
                _status.text("생성 완료!")
                st.success(f"A/B 버전 각 {_banner_n}장 생성 완료!")
            except Exception as _gen_err:
                st.error(f"생성 실패: {_gen_err}")

        # ── 결과 표시 ──
        if st.session_state.get("dp_banners_a") and st.session_state.get("dp_banners_b"):
            st.markdown("---")
            _tab_a, _tab_b = st.tabs(["📋 A버전 (심플형)", "🎨 B버전 (감성형)"])
            _copy_list_r = st.session_state.get("dp_copy_list", [])
            _blabels = ["히어로", "특징", "사용법", "리뷰", "정보"]

            with _tab_a:
                for _i, _img in enumerate(st.session_state["dp_banners_a"]):
                    _lbl = _blabels[_i] if _i < len(_blabels) else f"배너{_i+1}"
                    st.markdown(f"**배너 {_i+1} — {_lbl}**")
                    st.image(_img, use_container_width=True)
                    if _i < len(_copy_list_r):
                        _cp = _copy_list_r[_i]
                        with st.expander(f"배너 {_i+1} 카피 확인/수정"):
                            st.text_area("헤드라인", value=_cp.get("headline", ""),
                                         key=f"dp_a_hl_{_i}", height=55)
                            st.text_area("서브텍스트", value=_cp.get("subtext", ""),
                                         key=f"dp_a_sub_{_i}", height=55)
                            st.text_area("본문", value=_cp.get("body", ""),
                                         key=f"dp_a_body_{_i}", height=70)

            with _tab_b:
                for _i, _img in enumerate(st.session_state["dp_banners_b"]):
                    _lbl = _blabels[_i] if _i < len(_blabels) else f"배너{_i+1}"
                    st.markdown(f"**배너 {_i+1} — {_lbl}**")
                    st.image(_img, use_container_width=True)
                    if _i < len(_copy_list_r):
                        _cp = _copy_list_r[_i]
                        with st.expander(f"배너 {_i+1} 카피 확인/수정"):
                            st.text_area("헤드라인", value=_cp.get("headline", ""),
                                         key=f"dp_b_hl_{_i}", height=55)
                            st.text_area("서브텍스트", value=_cp.get("subtext", ""),
                                         key=f"dp_b_sub_{_i}", height=55)
                            st.text_area("본문", value=_cp.get("body", ""),
                                         key=f"dp_b_body_{_i}", height=70)

            # ── ZIP 다운로드 ──
            st.markdown("---")
            if st.button("⬇️ A/B 버전 ZIP 다운로드 준비",
                         key="dp_zip_prep_btn", use_container_width=True):
                with st.spinner("ZIP 파일 생성 중..."):
                    try:
                        _safe_name = re.sub(r'[^\w가-힣]', '_', dp_product_name)[:20] or "상품"
                        _zip_buf = create_detail_zip(
                            st.session_state["dp_banners_a"],
                            st.session_state["dp_banners_b"],
                            _safe_name
                        )
                        st.download_button(
                            label="⬇️ A/B 버전 ZIP 다운로드",
                            data=_zip_buf,
                            file_name=f"{_safe_name}_상세페이지.zip",
                            mime="application/zip",
                            type="primary",
                            use_container_width=True,
                            key="dp_zip_dl_btn"
                        )
                    except Exception as _ze:
                        st.error(f"ZIP 생성 실패: {_ze}")

# ─────────────────────────────────────────────
# 📖 바이럴 백과사전 메뉴
# ─────────────────────────────────────────────
elif menu == "📖 바이럴 백과사전":
    st.markdown("## 📖 바이럴 백과사전")
    st.caption("게시글을 관리하고 aligomedia.co.kr/blog/ 에 발행합니다.")

    _vb_tab1, _vb_tab2 = st.tabs(["📋 게시글 목록", "✍️ 새 글 작성"])

    # ── 탭1: 게시글 목록 ──
    with _vb_tab1:
        st.markdown("### 📋 발행된 게시글")

        # 연결 상태 진단
        with st.expander("🔍 연결 상태 확인"):
            _diag_tok = st.secrets.get("NETLIFY_TOKEN", "")
            _diag_sid = st.secrets.get("NETLIFY_SITE_ID", "")
            st.write(f"- NETLIFY_TOKEN: {'✅ 있음' if _diag_tok else '❌ 없음'}")
            st.write(f"- NETLIFY_SITE_ID: {'✅ 있음 (' + _diag_sid[:8] + '...)' if _diag_sid else '❌ 없음'}")
            _diag_ws = get_viral_sheet()
            st.write(f"- 구글시트 연결: {'✅ 성공 (바이럴백과사전 탭)' if _diag_ws else '❌ 실패'}")

        with st.spinner("게시글 불러오는 중..."):
            _vb_posts = get_viral_posts()

        if not _vb_posts:
            st.info("아직 게시된 글이 없습니다. '새 글 작성' 탭에서 첫 번째 글을 작성해보세요!")
        else:
            st.success(f"총 **{len(_vb_posts)}**개의 게시글이 있습니다.")
            for _vb_p in _vb_posts:
                _vc = get_view_count(_vb_p["slug"])
                with st.expander(f"📄 {_vb_p['제목']}  ·  {_vb_p['날짜'][:10]}  |  👁 {_vc:,}회"):
                    st.caption(f"🔑 slug: `{_vb_p['slug']}`")
                    st.caption(f"👁 조회수: **{_vc:,}회** (5분 캐시)")
                    if _vb_p["해시태그"]:
                        st.caption(f"🏷️ 태그: {_vb_p['해시태그']}")
                    if _vb_p["요약"]:
                        st.caption(f"📝 요약: {_vb_p['요약']}")
                    st.markdown(
                        f"[🔗 게시글 보기 (새 창)](https://aligomedia.co.kr/blog/{_vb_p['slug']}/)",
                        unsafe_allow_html=True)
                    if st.button("🗑️ 삭제", key=f"vb_del_{_vb_p['slug']}", type="secondary"):
                        with st.spinner("삭제 중..."):
                            _del_ok = delete_viral_post(_vb_p["slug"])
                            if _del_ok:
                                _remaining = [x for x in _vb_posts if x["slug"] != _vb_p["slug"]]
                                _del_idx_html = generate_blog_index_html(_remaining).encode("utf-8")
                                _vtok = st.secrets.get("NETLIFY_TOKEN", "")
                                _vsid = st.secrets.get("NETLIFY_SITE_ID", "")
                                if _vtok and _vsid:
                                    deploy_blog_incremental(
                                        _vtok, _vsid,
                                        {"blog/index.html": _del_idx_html})
                                st.success("✅ 삭제 완료!")
                                st.rerun()
                            else:
                                st.error("삭제 실패. 시트 연결을 확인해주세요.")

    # ── 탭2: 새 글 작성 ──
    with _vb_tab2:
        st.markdown("### ✍️ 새 글 작성")

        # 기본 정보
        _vb_title = st.text_input(
            "📌 제목 *",
            placeholder="예: 쿠팡 리뷰 마케팅 완벽 가이드 2025",
            key="vb_title")
        _vb_tags = st.text_input(
            "🏷️ 해시태그 (쉼표로 구분)",
            placeholder="예: 쿠팡리뷰, 리뷰마케팅, 바이럴마케팅",
            key="vb_tags")
        _vb_summary = st.text_area(
            "📝 요약 (검색 결과에 표시됨) *",
            placeholder="이 글에서 다루는 내용을 2~3줄로 요약해주세요. 네이버 검색 결과에 노출됩니다.",
            height=90, key="vb_summary")

        st.markdown("---")
        st.markdown("**📄 본문 — 블록 에디터**")
        st.caption("텍스트와 이미지를 원하는 순서로 자유롭게 쌓아가세요.")

        # ── 블록 초기화 ──
        import time as _time
        if "vb_blocks" not in st.session_state:
            st.session_state["vb_blocks"] = [
                {"type": "text", "id": "blk_init", "url": None}
            ]

        _vb_blocks = st.session_state["vb_blocks"]
        _vb_to_delete = None
        _vb_move = None  # (idx, direction)

        # ── 블록 렌더링 ──
        try:
            from streamlit_quill import st_quill as _st_quill
            _quill_ok = True
        except ImportError:
            _quill_ok = False

        for _bi, _blk in enumerate(_vb_blocks):
            _bid = _blk["id"]
            # 컨트롤 행
            _bc1, _bc2, _bc3, _bc4 = st.columns([1, 1, 9, 1])
            with _bc1:
                if st.button("↑", key=f"vbup_{_bid}", disabled=(_bi == 0), help="위로"):
                    _vb_move = (_bi, "up")
            with _bc2:
                if st.button("↓", key=f"vbdn_{_bid}", disabled=(_bi == len(_vb_blocks) - 1), help="아래로"):
                    _vb_move = (_bi, "down")
            with _bc3:
                if _blk["type"] == "text":
                    st.caption(f"📝 텍스트 블록 {_bi + 1}")
                else:
                    st.caption(f"🖼️ 이미지 블록 {_bi + 1}")
            with _bc4:
                if st.button("✕", key=f"vbdel_{_bid}", help="삭제"):
                    _vb_to_delete = _bi

            # 블록 내용
            if _blk["type"] == "text":
                if _quill_ok:
                    _raw = _st_quill(
                        placeholder="텍스트 입력... (글자 크기·색상·정렬 모두 가능)",
                        html=True,
                        key=f"q_{_bid}"
                    )
                    # rerun 후에도 내용 보존
                    _skey = f"__vbtxt_{_bid}"
                    if _raw is not None:
                        _clean = (_raw or "").strip()
                        if _clean and _clean not in ("<p><br></p>", "<p></p>"):
                            st.session_state[_skey] = _raw
                else:
                    _raw = st.text_area(
                        "텍스트 입력", height=150,
                        key=f"ta_{_bid}",
                        value=st.session_state.get(f"__vbtxt_{_bid}", ""))
                    if _raw:
                        st.session_state[f"__vbtxt_{_bid}"] = _raw

            elif _blk["type"] == "image":
                _stored_url = _blk.get("url")
                if _stored_url:
                    st.image(_stored_url, width=420)
                    if st.button("🔄 이미지 변경", key=f"vbchg_{_bid}"):
                        _vb_blocks[_bi]["url"] = None
                        st.rerun()
                else:
                    _img_file = st.file_uploader(
                        "이미지 파일 선택 (jpg / png / gif / webp)",
                        type=["jpg", "jpeg", "png", "gif", "webp"],
                        key=f"imgup_{_bid}"
                    )
                    if _img_file:
                        _vbtok = st.secrets.get("NETLIFY_TOKEN", "")
                        _vbsid = st.secrets.get("NETLIFY_SITE_ID", "")
                        if not _vbtok:
                            st.error("❌ NETLIFY_TOKEN 시크릿이 없습니다. Streamlit Cloud 시크릿 설정을 확인하세요.")
                        elif not _vbsid:
                            st.error("❌ NETLIFY_SITE_ID 시크릿이 없습니다. Streamlit Cloud 시크릿 설정을 확인하세요.")
                        else:
                            with st.spinner(f"{_img_file.name} 업로드 중..."):
                                _iurl, _imsg = upload_blog_image(
                                    _vbtok, _vbsid, _img_file.read(), _img_file.name)
                            if _iurl:
                                _vb_blocks[_bi]["url"] = _iurl
                                st.image(_iurl, width=420)
                                st.success("✅ 업로드 완료!")
                                st.rerun()
                            else:
                                st.error(f"업로드 실패: {_imsg}")

            st.markdown('<div style="height:6px"></div>', unsafe_allow_html=True)

        # ── 블록 이동 / 삭제 처리 ──
        if _vb_move:
            _mi, _mdir = _vb_move
            if _mdir == "up" and _mi > 0:
                _vb_blocks[_mi], _vb_blocks[_mi - 1] = _vb_blocks[_mi - 1], _vb_blocks[_mi]
            elif _mdir == "down" and _mi < len(_vb_blocks) - 1:
                _vb_blocks[_mi], _vb_blocks[_mi + 1] = _vb_blocks[_mi + 1], _vb_blocks[_mi]
            st.rerun()

        if _vb_to_delete is not None:
            _vb_blocks.pop(_vb_to_delete)
            st.rerun()

        # ── 블록 추가 버튼 ──
        _badd1, _badd2 = st.columns(2)
        with _badd1:
            if st.button("➕ 텍스트 블록 추가", use_container_width=True, key="vb_add_text"):
                _vb_blocks.append({"type": "text", "id": f"t{int(_time.time()*1000)}", "url": None})
                st.rerun()
        with _badd2:
            if st.button("🖼️ 이미지 블록 추가", use_container_width=True, key="vb_add_img"):
                _vb_blocks.append({"type": "image", "id": f"i{int(_time.time()*1000)}", "url": None})
                st.rerun()

        st.markdown("---")

        # ── 발행 버튼 ──
        _vb_col1, _vb_col2 = st.columns([2, 1])
        with _vb_col1:
            _vb_publish = st.button(
                "🚀 발행하기",
                key="vb_publish_btn",
                type="primary",
                use_container_width=True)
        with _vb_col2:
            st.caption("발행 시 홈페이지에 즉시 반영됩니다.")

        if _vb_publish:
            _vb_body_final = blocks_to_html(_vb_blocks)
            _vb_err = []
            if not (_vb_title or "").strip():
                _vb_err.append("제목을 입력해주세요.")
            if not (_vb_summary or "").strip():
                _vb_err.append("요약을 입력해주세요.")
            if not _vb_body_final.strip():
                _vb_err.append("본문 블록에 내용을 입력해주세요.")
            if _vb_err:
                for _e in _vb_err:
                    st.error(_e)
            else:
                _vtok2 = st.secrets.get("NETLIFY_TOKEN", "")
                _vsid2 = st.secrets.get("NETLIFY_SITE_ID", "")

                with st.spinner("게시글 저장 및 배포 중..."):
                    _vb_slug = make_slug(_vb_title.strip())
                    _vb_post_data = {
                        "날짜": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "slug": _vb_slug,
                        "제목": _vb_title.strip(),
                        "해시태그": (_vb_tags or "").strip(),
                        "요약": _vb_summary.strip(),
                        "본문HTML": _vb_body_final,
                    }

                    _vb_saved = save_viral_post(
                        _vb_slug,
                        _vb_title.strip(),
                        (_vb_tags or "").strip(),
                        _vb_summary.strip(),
                        _vb_body_final
                    )

                    if _vb_saved:
                        _vb_post_html = generate_post_html(_vb_post_data).encode("utf-8")
                        _vb_all_posts = get_viral_posts()
                        _vb_idx_html = generate_blog_index_html(_vb_all_posts).encode("utf-8")

                        if _vtok2 and _vsid2:
                            _vb_ok, _vb_msg = deploy_blog_incremental(
                                _vtok2, _vsid2,
                                {
                                    f"blog/{_vb_slug}/index.html": _vb_post_html,
                                    "blog/index.html": _vb_idx_html,
                                }
                            )
                            if _vb_ok:
                                # 발행 후 블록 초기화
                                st.session_state["vb_blocks"] = [
                                    {"type": "text", "id": "blk_init", "url": None}]
                                st.success("✅ 발행 완료!")
                                st.markdown(
                                    f"🔗 **게시글 주소:** https://aligomedia.co.kr/blog/{_vb_slug}/")
                                st.balloons()
                            else:
                                st.error(f"Netlify 배포 실패: {_vb_msg}")
                                st.info("구글시트에는 저장되었습니다.")
                        else:
                            st.warning("Netlify 시크릿 설정이 없어 홈페이지 배포는 건너뜁니다.\n구글시트에는 저장되었습니다.")
                    else:
                        st.error("구글시트 저장 실패. 시트 연결을 확인해주세요.")
